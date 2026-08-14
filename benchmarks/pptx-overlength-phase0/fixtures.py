#!/usr/bin/env python3

import argparse
import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


TARGET_NAME = "SPARKCLAW_PHASE0_TARGET"
EMU_PER_INCH = 914400
POINTS_PER_INCH = 72.0


CASES = (
    {
        "id": "latin_fit",
        "text": "Quarterly revenue improved eighteen percent.",
        "expect": "accept",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 7.5, 1.0),
    },
    {
        "id": "cjk_fit",
        "text": "季度收入同比增长百分之十八，核心业务保持稳定。",
        "expect": "accept",
        "font": "Noto Sans CJK SC",
        "box": (1.2, 1.5, 7.5, 1.0),
    },
    {
        "id": "mixed_fit",
        "text": "FY2026 收入增长 18%，retention 保持稳定。",
        "expect": "accept",
        "font": "Noto Sans CJK SC",
        "box": (1.2, 1.5, 7.5, 1.0),
    },
    {
        "id": "latin_clipped",
        "text": "LATIN-TAIL This deliberately oversized sentence must retain every final token including END-LATIN-9472.",
        "expect": "reject",
        "reason": "clipped_text",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 2.4, 0.38),
    },
    {
        "id": "cjk_clipped",
        "text": "中文超长内容必须完整呈现并保留最后的唯一标记终点九四七二否则必须拒绝",
        "expect": "reject",
        "reason": "clipped_text",
        "font": "Noto Sans CJK SC",
        "box": (1.2, 1.5, 2.1, 0.38),
    },
    {
        "id": "mixed_clipped",
        "text": "MIXED-START 中英文混合的超长内容 must preserve the unique suffix 终点-END-5831.",
        "expect": "reject",
        "reason": "clipped_text",
        "font": "Noto Sans CJK SC",
        "box": (1.2, 1.5, 2.3, 0.38),
    },
    {
        "id": "duplicate_attribution",
        "text": "DUPLICATE-ATTRIBUTION-7319",
        "expect": "reject",
        "reason": "ambiguous_attribution",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 5.0, 0.7),
        "duplicate": True,
    },
    {
        "id": "fully_occluded",
        "text": "OCCLUDED-TARGET-4826",
        "expect": "reject",
        "reason": "occluded",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 5.0, 0.7),
        "overlay": "opaque",
    },
    {
        "id": "partly_occluded",
        "text": "PARTIAL-OCCLUSION-2648",
        "expect": "reject",
        "reason": "occluded",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 5.0, 0.7),
        "overlay": "partial",
    },
    {
        "id": "transparent_overlay",
        "text": "TRANSPARENCY-TARGET-8135",
        "expect": "reject",
        "reason": "occluded",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 5.0, 0.7),
        "overlay": "transparent",
    },
    {
        "id": "same_color_concealment",
        "text": "SAME-COLOR-TARGET-3094",
        "expect": "reject",
        "reason": "same_color_concealment",
        "font": "Liberation Sans",
        "font_color": "FFFFFF",
        "box": (1.2, 1.5, 5.0, 0.7),
    },
    {
        "id": "soft_break",
        "text": "First verified line\nSecond verified line",
        "expect": "accept",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 6.0, 1.3),
        "soft_break": True,
    },
    {
        "id": "bullet",
        "text": "Evidence-bound bullet content",
        "expect": "accept",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 6.0, 0.8),
        "bullet": True,
    },
    {
        "id": "missing_font",
        "text": "MISSING-FONT-TARGET-6621",
        "expect": "reject",
        "reason": "font_unavailable",
        "font": "SparkClaw Missing Font 2026",
        "box": (1.2, 1.5, 5.0, 0.7),
    },
    {
        "id": "rotated_target",
        "text": "ROTATED-UNSUPPORTED-1753",
        "expect": "reject",
        "reason": "unsupported_target",
        "font": "Liberation Sans",
        "box": (1.2, 1.5, 5.0, 0.7),
        "rotation": 12.0,
        "supported": False,
    },
)


FOUR_BY_THREE_CASES = (
    {
        "id": "aspect_4_3_fit",
        "text": "A four by three presentation remains inside its target.",
        "expect": "accept",
        "font": "Liberation Sans",
        "box": (1.0, 1.3, 6.5, 0.9),
    },
    {
        "id": "aspect_4_3_clipped",
        "text": "FOUR-BY-THREE clipped content must not lose the final marker END-4438.",
        "expect": "reject",
        "reason": "clipped_text",
        "font": "Liberation Sans",
        "box": (1.0, 1.3, 2.0, 0.38),
    },
)


def _hex_color(value):
    return RGBColor.from_string(value)


def _set_fill_transparency(shape, percentage):
    solid = shape.fill._xPr.solidFill
    color = solid.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    if color is None:
        return
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int((100 - percentage) * 1000)))
    color.append(alpha)


def _set_bullet(paragraph):
    properties = paragraph._p.get_or_add_pPr()
    for child in list(properties):
        if child.tag.endswith("}buNone") or child.tag.endswith("}buChar"):
            properties.remove(child)
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "\u2022")
    properties.append(bullet)


def _add_target(slide, case):
    left, top, width, height = case["box"]
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = TARGET_NAME
    shape.rotation = float(case.get("rotation", 0.0))
    frame = shape.text_frame
    frame.clear()
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    text = case["text"]
    if case.get("soft_break"):
        first, second = text.split("\n", 1)
        run = paragraph.add_run()
        run.text = first
        paragraph.add_line_break()
        second_run = paragraph.add_run()
        second_run.text = second
        runs = (run, second_run)
    else:
        run = paragraph.add_run()
        run.text = text
        runs = (run,)
    if case.get("bullet"):
        _set_bullet(paragraph)
    for current in runs:
        current.font.name = case["font"]
        current.font.size = Pt(22)
        current.font.color.rgb = _hex_color(case.get("font_color", "111111"))
    return shape


def _add_overlay(slide, case):
    mode = case.get("overlay")
    if not mode:
        return
    left, top, width, height = case["box"]
    if mode == "partial":
        left += width * 0.48
        width *= 0.52
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    overlay.name = "SPARKCLAW_PHASE0_OCCLUDER"
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = _hex_color("FFFFFF")
    overlay.line.fill.background()
    if mode == "transparent":
        _set_fill_transparency(overlay, 50)


def _add_case_slide(prs, case):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _hex_color("FFFFFF")
    target = _add_target(slide, case)
    if case.get("duplicate"):
        duplicate = copy.deepcopy(target.element)
        slide.shapes._spTree.insert_element_before(duplicate, "p:extLst")
        duplicate_shape = slide.shapes[-1]
        duplicate_shape.left = Inches(1.2)
        duplicate_shape.top = Inches(3.2)
        duplicate_shape.name = "SPARKCLAW_PHASE0_DUPLICATE"
    _add_overlay(slide, case)
    marker = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(7.0), Inches(0.3))
    marker.text_frame.paragraphs[0].text = "fixture:" + case["id"]
    marker.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    marker.text_frame.paragraphs[0].runs[0].font.color.rgb = _hex_color("777777")
    return target


def _target_region(shape):
    frame = shape.text_frame
    left = (shape.left + frame.margin_left) / EMU_PER_INCH * POINTS_PER_INCH
    top = (shape.top + frame.margin_top) / EMU_PER_INCH * POINTS_PER_INCH
    width = (shape.width - frame.margin_left - frame.margin_right) / EMU_PER_INCH * POINTS_PER_INCH
    height = (shape.height - frame.margin_top - frame.margin_bottom) / EMU_PER_INCH * POINTS_PER_INCH
    return {
        "x": round(left, 3),
        "y": round(top, 3),
        "width": round(width, 3),
        "height": round(height, 3),
    }


def _new_presentation(aspect):
    prs = Presentation()
    while prs.slides:
        slide_id = prs.slides._sldIdLst[0]
        rel_id = slide_id.rId
        prs.slides._sldIdLst.remove(slide_id)
        prs.part.drop_rel(rel_id)
    if aspect == "16:9":
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    return prs


def _save_variants(prs, output_dir, stem):
    candidate_path = output_dir / (stem + ".pptx")
    prs.save(candidate_path)

    without = Presentation(candidate_path)
    for slide in without.slides:
        target = next(shape for shape in slide.shapes if shape.name == TARGET_NAME)
        for paragraph in target.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""
    without_path = output_dir / (stem + "-without-target.pptx")
    without.save(without_path)

    topmost = Presentation(candidate_path)
    for slide in topmost.slides:
        target = next(shape for shape in slide.shapes if shape.name == TARGET_NAME)
        parent = target.element.getparent()
        parent.remove(target.element)
        parent.insert_element_before(target.element, "p:extLst")
    topmost_path = output_dir / (stem + "-target-topmost.pptx")
    topmost.save(topmost_path)
    return candidate_path, without_path, topmost_path


def generate(output_dir):
    output_dir.mkdir(parents=True, exist_ok=True)
    manifest = {
        "schema_version": "sparkclaw.pptx_overlength.fixture_manifest.v1",
        "decks": [],
    }
    for stem, aspect, cases in (
        ("phase0-16x9", "16:9", CASES),
        ("phase0-4x3", "4:3", FOUR_BY_THREE_CASES),
    ):
        prs = _new_presentation(aspect)
        records = []
        for page_index, case in enumerate(cases, start=1):
            target = _add_case_slide(prs, case)
            records.append({
                "id": case["id"],
                "page": page_index,
                "text": case["text"],
                "expect": case["expect"],
                "reason": case.get("reason", ""),
                "font": case["font"],
                "supported": case.get("supported", True),
                "region_pt": _target_region(target),
            })
        candidate, without, topmost = _save_variants(prs, output_dir, stem)
        manifest["decks"].append({
            "id": stem,
            "aspect": aspect,
            "candidate": candidate.name,
            "without_target": without.name,
            "target_topmost": topmost.name,
            "cases": records,
        })
    manifest_path = output_dir / "fixtures.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return manifest_path


def main():
    parser = argparse.ArgumentParser(description="Generate deterministic PPTX Phase 0 fixtures.")
    parser.add_argument("output_dir", type=Path)
    args = parser.parse_args()
    print(generate(args.output_dir))


if __name__ == "__main__":
    main()
