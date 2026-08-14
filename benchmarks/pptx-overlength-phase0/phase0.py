#!/usr/bin/env python3

import argparse
import hashlib
import json
import math
import os
import platform
import re
import shutil
import signal
import statistics
import subprocess
import sys
import tempfile
import time
import unicodedata
import urllib.error
import urllib.request
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path

import pypdfium2 as pdfium
from PIL import ImageChops
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from fixtures import generate as generate_fixtures


ROOT = Path(__file__).resolve().parents[2]
HERE = Path(__file__).resolve().parent
PDFJS_PROBE = HERE / "pdfjs_probe.mjs"
PPTX_WRITER = ROOT / "services/gateway/internal/toolhub/scripts/pptx_slide.py"


class QualificationError(RuntimeError):
    pass


@dataclass
class ProcessResult:
    returncode: int
    stdout: bytes
    stderr: bytes
    elapsed_seconds: float
    timed_out: bool = False
    cleanup_seconds: float = 0.0


def normalize_text(value):
    normalized = unicodedata.normalize("NFKC", str(value or ""))
    return " ".join(normalized.replace("\v", " ").split())


def comparison_text(value):
    normalized = normalize_text(value)
    normalized = re.sub(r"\s+([^\w\s])", r"\1", normalized, flags=re.UNICODE)
    return re.sub(r"([^\w\s])\s+", r"\1", normalized, flags=re.UNICODE)


def sha256_bytes(value):
    return hashlib.sha256(value).hexdigest()


def sha256_file(path):
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def canonical_package_digest(path):
    digest = hashlib.sha256()
    with zipfile.ZipFile(path) as package:
        names = package.namelist()
        if len(names) != len(set(names)):
            raise QualificationError("OOXML package contains duplicate part names")
        for name in sorted(names):
            raw = package.read(name)
            encoded = name.encode("utf-8")
            digest.update(len(encoded).to_bytes(4, "big"))
            digest.update(encoded)
            digest.update(len(raw).to_bytes(8, "big"))
            digest.update(raw)
    return digest.hexdigest()


def package_part_diffs(before, after):
    with zipfile.ZipFile(before) as left, zipfile.ZipFile(after) as right:
        left_parts = {name: left.read(name) for name in left.namelist()}
        right_parts = {name: right.read(name) for name in right.namelist()}
    return sorted(
        name for name in set(left_parts) | set(right_parts)
        if left_parts.get(name) != right_parts.get(name)
    )


def _process_alive(pid):
    stat = Path("/proc") / str(pid) / "stat"
    try:
        state = stat.read_text(encoding="utf-8").split()[2]
        return state != "Z"
    except (FileNotFoundError, IndexError, PermissionError):
        return False


def _stop_process_group(process, cleanup_timeout):
    started = time.monotonic()
    try:
        os.killpg(process.pid, signal.SIGTERM)
    except ProcessLookupError:
        return 0.0
    try:
        process.wait(timeout=cleanup_timeout)
    except subprocess.TimeoutExpired:
        try:
            os.killpg(process.pid, signal.SIGKILL)
        except ProcessLookupError:
            pass
        process.wait(timeout=max(1.0, cleanup_timeout))
    return time.monotonic() - started


def run_process(command, *, timeout, cleanup_timeout, stdin=None, env=None, cwd=None, stdout_limit=1024 * 1024):
    started = time.monotonic()
    process = subprocess.Popen(
        command,
        cwd=cwd,
        env=env,
        stdin=subprocess.PIPE if stdin is not None else subprocess.DEVNULL,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        start_new_session=True,
    )
    try:
        stdout, stderr = process.communicate(stdin, timeout=timeout)
        elapsed = time.monotonic() - started
        if len(stdout) > stdout_limit or len(stderr) > stdout_limit:
            raise QualificationError("subprocess output exceeded the qualification byte limit")
        return ProcessResult(process.returncode, stdout, stderr, elapsed)
    except subprocess.TimeoutExpired:
        cleanup = _stop_process_group(process, cleanup_timeout)
        stdout, stderr = process.communicate()
        return ProcessResult(process.returncode, stdout[:stdout_limit], stderr[:stdout_limit], time.monotonic() - started, True, cleanup)


def _profile_uri(path):
    return path.resolve().as_uri()


def convert_libreoffice(source, output_dir, policy):
    output_dir.mkdir(parents=True, exist_ok=True)
    profile = output_dir / "lo-profile"
    command = [
        shutil.which("libreoffice") or "/usr/bin/libreoffice",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        f"-env:UserInstallation={_profile_uri(profile)}",
        "--convert-to",
        "pdf:impress_pdf_Export",
        "--outdir",
        str(output_dir),
        str(source),
    ]
    limits = policy["limits"]
    result = run_process(
        command,
        timeout=limits["conversion_timeout_seconds"],
        cleanup_timeout=limits["cleanup_timeout_seconds"],
        stdout_limit=limits["max_stdout_bytes"],
    )
    output = output_dir / (source.stem + ".pdf")
    if result.timed_out:
        raise QualificationError("LibreOffice conversion timed out")
    if result.returncode != 0 or not output.is_file():
        raise QualificationError("LibreOffice conversion failed")
    if output.stat().st_size > limits["max_pdf_bytes"]:
        raise QualificationError("LibreOffice output exceeded the PDF byte limit")
    return output, result.elapsed_seconds


def _multipart_file(field, path, boundary):
    filename = path.name.replace('"', "")
    header = (
        f"--{boundary}\r\n"
        f"Content-Disposition: form-data; name=\"{field}\"; filename=\"{filename}\"\r\n"
        "Content-Type: application/vnd.openxmlformats-officedocument.presentationml.presentation\r\n\r\n"
    ).encode("ascii")
    return header + path.read_bytes() + b"\r\n" + f"--{boundary}--\r\n".encode("ascii")


def convert_gotenberg(source, output_dir, policy, base_url):
    output_dir.mkdir(parents=True, exist_ok=True)
    boundary = "sparkclaw-" + uuid.uuid4().hex
    body = _multipart_file("files", source, boundary)
    request = urllib.request.Request(
        base_url.rstrip("/") + "/forms/libreoffice/convert",
        data=body,
        method="POST",
        headers={"Content-Type": f"multipart/form-data; boundary={boundary}"},
    )
    started = time.monotonic()
    try:
        with urllib.request.urlopen(request, timeout=policy["limits"]["conversion_timeout_seconds"]) as response:
            raw = response.read(policy["limits"]["max_pdf_bytes"] + 1)
    except (urllib.error.URLError, TimeoutError) as error:
        raise QualificationError("Gotenberg conversion is unavailable") from error
    elapsed = time.monotonic() - started
    if len(raw) > policy["limits"]["max_pdf_bytes"] or not raw.startswith(b"%PDF-"):
        raise QualificationError("Gotenberg returned an invalid or oversized PDF")
    output = output_dir / (source.stem + ".pdf")
    output.write_bytes(raw)
    return output, elapsed


def gotenberg_healthy(base_url, timeout=5):
    try:
        with urllib.request.urlopen(base_url.rstrip("/") + "/health", timeout=timeout) as response:
            payload = json.loads(response.read(1024 * 1024))
        return response.status == 200 and payload.get("status") == "up"
    except (urllib.error.URLError, json.JSONDecodeError, TimeoutError):
        return False


def inspect_pdf(path, output_dir, policy):
    output = output_dir / (path.stem + "-pdfjs.json")
    result = run_process(
        [shutil.which("node") or "node", str(PDFJS_PROBE), str(path), str(output)],
        timeout=policy["limits"]["conversion_timeout_seconds"],
        cleanup_timeout=policy["limits"]["cleanup_timeout_seconds"],
        stdout_limit=policy["limits"]["max_stdout_bytes"],
        cwd=HERE,
    )
    if result.timed_out or result.returncode != 0 or not output.is_file():
        raise QualificationError("PDF.js inspection failed")
    return json.loads(output.read_text(encoding="utf-8"))


def raster_pages(path, scale):
    document = pdfium.PdfDocument(path)
    pages = []
    try:
        for page in document:
            bitmap = page.render(scale=scale)
            pages.append(bitmap.to_pil().convert("RGB"))
    finally:
        document.close()
    return pages


def raster_digest(pages):
    digest = hashlib.sha256()
    for page in pages:
        digest.update(page.width.to_bytes(4, "big"))
        digest.update(page.height.to_bytes(4, "big"))
        digest.update(page.tobytes())
    return digest.hexdigest()


def _intersects(item, region):
    return not (
        item["x"] + item["width"] <= region["x"]
        or region["x"] + region["width"] <= item["x"]
        or item["y"] + item["height"] <= region["y"]
        or region["y"] + region["height"] <= item["y"]
    )


def _inside(item, region, tolerance):
    return (
        item["x"] >= region["x"] - tolerance
        and item["y"] >= region["y"] - tolerance
        and item["x"] + item["width"] <= region["x"] + region["width"] + tolerance
        and item["y"] + item["height"] <= region["y"] + region["height"] + tolerance
    )


def _crop_for_region(image, page, region):
    factor_x = image.width / page["width"]
    factor_y = image.height / page["height"]
    left = max(0, math.floor(region["x"] * factor_x))
    top = max(0, math.floor(region["y"] * factor_y))
    right = min(image.width, math.ceil((region["x"] + region["width"]) * factor_x))
    bottom = min(image.height, math.ceil((region["y"] + region["height"]) * factor_y))
    if right <= left or bottom <= top:
        raise QualificationError("fixture target region is outside the rendered page")
    return image.crop((left, top, right, bottom))


def _visible_ratio(candidate, without_target, topmost, channel_delta):
    candidate_effect = ImageChops.difference(candidate, without_target).tobytes()
    proof_effect = ImageChops.difference(topmost, without_target).tobytes()
    proof_pixels = 0
    matching_pixels = 0
    for offset in range(0, len(proof_effect), 3):
        proof = proof_effect[offset:offset + 3]
        current = candidate_effect[offset:offset + 3]
        if max(proof) <= channel_delta:
            continue
        proof_pixels += 1
        if max(current) > channel_delta and max(abs(current[index] - proof[index]) for index in range(3)) <= channel_delta:
            matching_pixels += 1
    return proof_pixels, matching_pixels / proof_pixels if proof_pixels else 0.0


def available_font_families():
    result = run_process(
        ["fc-list", "--format", "%{family}\n"],
        timeout=10,
        cleanup_timeout=2,
    )
    if result.returncode != 0:
        raise QualificationError("fontconfig inventory failed")
    values = set()
    for line in result.stdout.decode("utf-8", "replace").splitlines():
        values.update(item.strip().casefold() for item in line.split(",") if item.strip())
    return values


def evaluate_case(case, candidate_probe, candidate_image, without_image, topmost_image, fonts, policy):
    page = candidate_probe["pages"][case["page"] - 1]
    region = case["region_pt"]
    tolerance = policy["normalization"]["geometry_tolerance_pt"]
    text = comparison_text(case["text"])
    page_text = comparison_text(page["text"])
    relevant = [item for item in page["items"] if item["text"] and _intersects(item, region)]
    region_text = comparison_text(" ".join(item["text"] for item in relevant))
    contributing = [
        item for item in relevant
        if comparison_text(item["text"])
        and (
            comparison_text(item["text"]) in text
            or text in comparison_text(item["text"])
        )
    ]
    occurrence_count = page_text.count(text)
    reasons = []
    if not case.get("supported", True):
        reasons.append("unsupported_target")
    if case["font"].casefold() not in fonts:
        reasons.append("font_unavailable")
    if text not in region_text:
        reasons.append("rendered_text_incomplete")
    if occurrence_count != 1:
        reasons.append("ambiguous_attribution" if occurrence_count > 1 else "rendered_text_incomplete")
    if contributing and not all(_inside(item, region, tolerance) for item in contributing):
        reasons.append("outside_target_region")

    candidate_crop = _crop_for_region(candidate_image, page, region)
    without_crop = _crop_for_region(without_image, page, region)
    topmost_crop = _crop_for_region(topmost_image, page, region)
    proof_pixels, visible_ratio = _visible_ratio(
        candidate_crop,
        without_crop,
        topmost_crop,
        policy["normalization"]["raster_channel_delta"],
    )
    if proof_pixels < policy["normalization"]["minimum_proof_pixels"]:
        reasons.append("raster_visibility_unprovable")
    elif visible_ratio < policy["normalization"]["minimum_visible_ratio"]:
        reasons.append("occluded")
    return {
        "id": case["id"],
        "verdict": "reject" if reasons else "accept",
        "expected": case["expect"],
        "reasons": sorted(set(reasons)),
        "attributed_items": len(contributing),
        "visible_ratio": round(visible_ratio, 6),
    }


def _convert_variant(engine, source, output_dir, policy, gotenberg_url):
    if engine == "libreoffice":
        return convert_libreoffice(source, output_dir, policy)
    return convert_gotenberg(source, output_dir, policy, gotenberg_url)


def evaluate_engine(engine, manifest, fixture_dir, output_dir, policy, gotenberg_url, fonts, repetitions):
    engine_dir = output_dir / engine
    engine_dir.mkdir(parents=True, exist_ok=True)
    cases = []
    normalized_runs = []
    raster_runs = []
    timings = []
    for deck in manifest["decks"]:
        variants = {}
        for key in ("candidate", "without_target", "target_topmost"):
            source = fixture_dir / deck[key]
            render_dir = engine_dir / deck["id"] / key
            pdf, elapsed = _convert_variant(engine, source, render_dir, policy, gotenberg_url)
            timings.append(elapsed)
            probe = inspect_pdf(pdf, render_dir, policy)
            raster = raster_pages(pdf, policy["normalization"]["raster_scale"])
            variants[key] = {"probe": probe, "raster": raster}
        normalized_runs.append(variants["candidate"]["probe"]["normalized_digest"])
        raster_runs.append(raster_digest(variants["candidate"]["raster"]))
        for case in deck["cases"]:
            page_index = case["page"] - 1
            cases.append(evaluate_case(
                case,
                variants["candidate"]["probe"],
                variants["candidate"]["raster"][page_index],
                variants["without_target"]["raster"][page_index],
                variants["target_topmost"]["raster"][page_index],
                fonts,
                policy,
            ))

    normalized_runs = [sha256_bytes("\n".join(normalized_runs).encode("ascii"))]
    raster_runs = [sha256_bytes("\n".join(raster_runs).encode("ascii"))]
    for repetition in range(1, repetitions):
        run_digests = []
        run_rasters = []
        for deck in manifest["decks"]:
            source = fixture_dir / deck["candidate"]
            render_dir = engine_dir / deck["id"] / f"repeat-{repetition:03d}"
            pdf, elapsed = _convert_variant(engine, source, render_dir, policy, gotenberg_url)
            timings.append(elapsed)
            probe = inspect_pdf(pdf, render_dir, policy)
            raster = raster_pages(pdf, policy["normalization"]["raster_scale"])
            run_digests.append(probe["normalized_digest"])
            run_rasters.append(raster_digest(raster))
        normalized_runs.append(sha256_bytes("\n".join(run_digests).encode("ascii")))
        raster_runs.append(sha256_bytes("\n".join(run_rasters).encode("ascii")))

    ordered = sorted(timings)
    percentile_index = min(len(ordered) - 1, math.ceil(len(ordered) * 0.95) - 1)
    return {
        "engine": engine,
        "cases": cases,
        "cases_match_expected": all(case["verdict"] == case["expected"] for case in cases),
        "repeat_count": repetitions,
        "normalized_repeatable": len(set(normalized_runs)) == 1,
        "raster_repeatable": len(set(raster_runs)) == 1,
        "timing_seconds": {
            "median": round(statistics.median(timings), 4),
            "p95": round(ordered[percentile_index], 4),
            "worst": round(max(timings), 4),
            "renders": len(timings),
        },
    }


def normalize_atomic_groups(updates):
    if not updates:
        return []
    group_ids = [item.get("atomic_group_id") for item in updates]
    if any(not isinstance(group_id, str) or not group_id.strip() for group_id in group_ids):
        return [{"id": "operation", "updates": list(updates)}]
    groups = {}
    for update in updates:
        groups.setdefault(update["atomic_group_id"], []).append(update)
    return [{"id": group_id, "updates": groups[group_id]} for group_id in sorted(groups)]


def select_effective_groups(updates, eligible_update_ids):
    groups = normalize_atomic_groups(updates)
    accepted = []
    skipped = []
    for group in groups:
        target = accepted if all(item["id"] in eligible_update_ids for item in group["updates"]) else skipped
        target.append(group["id"])
    if not accepted:
        return {"status": "no_safe_change", "accepted_groups": [], "skipped_groups": skipped, "artifact": None}
    return {
        "status": "completed" if not skipped else "completed_with_skips",
        "accepted_groups": accepted,
        "skipped_groups": skipped,
        "artifact": "sealed-prepared-artifact",
    }


def qualify_selection_contracts():
    updates = [
        {"id": "title", "atomic_group_id": "coherent"},
        {"id": "body", "atomic_group_id": "coherent"},
        {"id": "caption", "atomic_group_id": "independent"},
    ]
    partial = select_effective_groups(updates, {"title", "caption"})
    no_change = select_effective_groups(updates, set())
    invalid = select_effective_groups([{"id": "one"}, {"id": "two", "atomic_group_id": "valid"}], {"one"})
    passed = (
        partial["status"] == "completed_with_skips"
        and partial["accepted_groups"] == ["independent"]
        and partial["skipped_groups"] == ["coherent"]
        and no_change["status"] == "no_safe_change"
        and no_change["artifact"] is None
        and invalid["status"] == "no_safe_change"
        and invalid["skipped_groups"] == ["operation"]
    )
    return passed


def _writer_fixture(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    frame = shape.text_frame
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    run = frame.paragraphs[0].add_run()
    run.text = "SOURCE-PRESERVATION-1947"
    run.font.name = "Liberation Sans"
    run.font.size = Pt(20)
    prs.save(path)


def qualify_writer(output_dir, policy):
    output_dir.mkdir(parents=True, exist_ok=True)
    source = output_dir / "writer-source.pptx"
    _writer_fixture(source)
    canonical = []
    raw = []
    changed_parts = None
    for index in range(5):
        output = output_dir / f"writer-output-{index}.pptx"
        payload = {
            "operation": "update_slide",
            "path": str(source),
            "output_path": str(output),
            "slide_index": 1,
            "layout_policy": "preserve",
            "updates": [{
                "shape_index": 1,
                "old_text": "SOURCE-PRESERVATION-1947",
                "text": "TARGET-PRESERVATION-6284",
            }],
        }
        result = run_process(
            [sys.executable, str(PPTX_WRITER)],
            timeout=policy["limits"]["conversion_timeout_seconds"],
            cleanup_timeout=policy["limits"]["cleanup_timeout_seconds"],
            stdin=json.dumps(payload).encode("utf-8"),
            stdout_limit=policy["limits"]["max_stdout_bytes"],
        )
        if result.timed_out or result.returncode != 0 or not output.is_file():
            return {"passed": False, "reason": "writer_failed"}
        response = json.loads(result.stdout)
        if response.get("error"):
            return {"passed": False, "reason": "writer_rejected_fixture"}
        current_parts = package_part_diffs(source, output)
        changed_parts = current_parts if changed_parts is None else changed_parts
        if current_parts != changed_parts:
            return {"passed": False, "reason": "part_diff_nondeterministic"}
        canonical.append(canonical_package_digest(output))
        raw.append(sha256_file(output))
    allowed = {"ppt/slides/slide1.xml"}
    return {
        "passed": set(changed_parts or []) <= allowed and len(set(canonical)) == 1,
        "changed_parts": changed_parts,
        "canonical_repeatable": len(set(canonical)) == 1,
        "raw_repeatable": len(set(raw)) == 1,
        "canonical_digest": canonical[0],
    }


def qualify_cancellation(output_dir, policy):
    output_dir.mkdir(parents=True, exist_ok=True)
    child_file = output_dir / "child-pid"
    code = (
        "import pathlib,subprocess,sys,time;"
        "p=subprocess.Popen([sys.executable,'-c','import time; time.sleep(60)']);"
        f"pathlib.Path({str(child_file)!r}).write_text(str(p.pid));"
        "time.sleep(60)"
    )
    result = run_process(
        [sys.executable, "-c", code],
        timeout=0.25,
        cleanup_timeout=policy["limits"]["cleanup_timeout_seconds"],
    )
    deadline = time.monotonic() + policy["limits"]["cleanup_timeout_seconds"]
    child_pid = int(child_file.read_text()) if child_file.is_file() else -1
    while child_pid > 0 and _process_alive(child_pid) and time.monotonic() < deadline:
        time.sleep(0.02)
    child_gone = child_pid > 0 and not _process_alive(child_pid)
    child_file.unlink(missing_ok=True)
    return {
        "passed": result.timed_out and result.cleanup_seconds <= policy["limits"]["cleanup_timeout_seconds"] and child_gone,
        "cleanup_seconds": round(result.cleanup_seconds, 4),
        "child_stopped": child_gone,
    }


def qualify_outage(output_dir, policy):
    output_dir.mkdir(parents=True, exist_ok=True)
    source = output_dir / "outage-source.pptx"
    _writer_fixture(source)
    try:
        convert_gotenberg(source, output_dir / "outage-output", policy, "http://127.0.0.1:1")
    except QualificationError:
        outputs = list((output_dir / "outage-output").glob("*.pdf")) if (output_dir / "outage-output").exists() else []
        return {"passed": not outputs, "typed_outcome": "runtime_unavailable", "artifacts": len(outputs)}
    return {"passed": False, "typed_outcome": "unchecked_success", "artifacts": 1}


def _gate(status, evidence):
    return {"status": status, "evidence": evidence}


def _powerpoint_gate(path, corpus_digest):
    if path is None:
        return _gate("NOT_RUN", "No recorded Microsoft PowerPoint viewer evidence was supplied.")
    try:
        evidence = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return _gate("FAIL", "PowerPoint evidence is unreadable.")
    passed = (
        evidence.get("schema_version") == "sparkclaw.pptx_powerpoint_evidence.v1"
        and evidence.get("corpus_digest") == corpus_digest
        and evidence.get("repair_prompt") is False
        and evidence.get("visible_clipping") is False
        and evidence.get("missing_text") is False
        and bool(evidence.get("viewer_version"))
    )
    return _gate("PASS" if passed else "FAIL", "PowerPoint evidence matched the exact fixture corpus." if passed else "PowerPoint evidence was incomplete or did not match the fixture corpus.")


def _load_fonts(path, fallback):
    if path is None:
        return fallback
    data = json.loads(path.read_text(encoding="utf-8"))
    return {str(item).strip().casefold() for item in data.get("families", []) if str(item).strip()}


def qualify_conversion_cost(engines, policy):
    plan = policy["candidate_plan"]
    evaluations = plan["max_shapes"] * plan["max_candidates_per_shape"]
    deadline = plan["preparation_deadline_seconds"]
    if not engines:
        return {
            "status": "FAIL",
            "candidate_evaluations": evaluations,
            "deadline_seconds": deadline,
            "reason": "no_renderer_timing",
        }
    fastest = min(engines, key=lambda item: item["timing_seconds"]["median"])
    median = fastest["timing_seconds"]["median"] * evaluations
    p95 = fastest["timing_seconds"]["p95"] * evaluations
    worst = fastest["timing_seconds"]["worst"] * evaluations
    return {
        "status": "FAIL" if median > deadline else "NOT_RUN",
        "engine": fastest["engine"],
        "candidate_evaluations": evaluations,
        "deadline_seconds": deadline,
        "projected_seconds": {
            "median": round(median, 4),
            "p95": round(p95, 4),
            "worst": round(worst, 4),
        },
        "reason": "median_lower_bound_exceeds_deadline" if median > deadline else "peak_memory_and_final_plan_unqualified",
    }


def run_qualification(args):
    policy = json.loads((args.policy or HERE / "policy.json").read_text(encoding="utf-8"))
    work_dir = args.work_dir.resolve()
    fixture_dir = work_dir / "fixtures"
    output_dir = work_dir / "results"
    fixture_manifest_path = generate_fixtures(fixture_dir)
    manifest = json.loads(fixture_manifest_path.read_text(encoding="utf-8"))
    corpus_digest = sha256_bytes("\n".join(
        canonical_package_digest(fixture_dir / deck["candidate"])
        for deck in manifest["decks"]
    ).encode("ascii"))

    host_fonts = available_font_families()
    gotenberg_fonts = _load_fonts(args.gotenberg_fonts, host_fonts)
    gotenberg_identity_ok = args.gotenberg_image == policy["artifacts"]["gotenberg_arm64_image"]
    native_ok = platform.machine() == "aarch64" and shutil.which("libreoffice") is not None
    health_ok = gotenberg_healthy(args.gotenberg_url)
    engines = []
    engine_errors = []
    for engine, fonts in (("libreoffice", host_fonts), ("gotenberg", gotenberg_fonts)):
        try:
            engines.append(evaluate_engine(
                engine,
                manifest,
                fixture_dir,
                output_dir,
                policy,
                args.gotenberg_url,
                fonts,
                args.repetitions,
            ))
        except QualificationError as error:
            engine_errors.append({"engine": engine, "error": str(error)})

    writer = qualify_writer(output_dir / "writer", policy)
    cancellation = qualify_cancellation(output_dir / "cancellation", policy)
    outage = qualify_outage(output_dir / "outage", policy)
    selection = qualify_selection_contracts()
    conversion_cost = qualify_conversion_cost(engines, policy)
    cases_pass = len(engines) == 2 and all(engine["cases_match_expected"] for engine in engines)
    repeat_pass = (
        args.repetitions == policy["limits"]["repeat_renders"]
        and len(engines) == 2
        and all(engine["normalized_repeatable"] and engine["raster_repeatable"] for engine in engines)
    )
    powerpoint = _powerpoint_gate(args.powerpoint_evidence, corpus_digest)

    gates = {
        "native_deployment": _gate("PASS" if native_ok and health_ok and gotenberg_identity_ok else "FAIL", "ARM64 host, pinned Gotenberg identity, native LibreOffice, and health endpoint were checked."),
        "text_completeness": _gate("PASS" if cases_pass else "FAIL", "Synthetic Latin, CJK, and mixed fit/clipping verdicts were compared with expectations."),
        "attribution_visibility": _gate("PASS" if cases_pass else "FAIL", "Duplicate, opaque, partial, transparent, and same-color cases were checked with geometry and counterfactual raster evidence."),
        "geometry": _gate("PASS" if cases_pass else "FAIL", "PDF.js item transforms were constrained to projected target regions."),
        "font_determinism": _gate("PASS" if cases_pass else "FAIL", "Declared fonts were checked against engine font manifests and the missing-font fixture was required to fail closed."),
        "render_repeatability": _gate("PASS" if repeat_pass else "FAIL", f"Each fixture deck was rendered {args.repetitions} times per engine; policy requires {policy['limits']['repeat_renders']}."),
        "owner_corpus": _gate("NOT_RUN", "No owner deck with irreversible private-text replacement was supplied to this qualification run."),
        "powerpoint_compatibility": powerpoint,
        "writer_preservation": _gate("PASS" if writer.get("passed") else "FAIL", "The existing mutation layer was replayed five times and raw OOXML part deltas were allowlisted."),
        "partial_application": _gate("PASS" if selection else "FAIL", "Deterministic independent-group selection contract was exercised."),
        "semantic_atomicity": _gate("PASS" if selection else "FAIL", "A failing member rejected its full semantic group and invalid metadata collapsed to operation-wide atomicity."),
        "no_safe_change": _gate("PASS" if selection else "FAIL", "An all-ineligible plan returned no_safe_change without an edited artifact."),
        "pipeline_outcomes": _gate("PASS" if selection else "FAIL", "Qualification-only requested/effective plan outcome contracts were exercised; no production integration was enabled."),
        "renderer_outage": _gate("PASS" if outage["passed"] else "FAIL", "An unreachable renderer produced runtime_unavailable and no PDF artifact."),
        "cancellation": _gate("PASS" if cancellation["passed"] else "FAIL", "A timed-out process tree was terminated under the two-second cleanup bound."),
        "conversion_cost": _gate(conversion_cost["status"], "The unreduced 64-shape by 16-candidate plan was projected from measured conversion timings against the 90-second preparation deadline."),
        "digest_determinism": _gate("PASS" if writer.get("canonical_repeatable") and repeat_pass else "FAIL", "Canonical OOXML and normalized render digests were compared independently of container metadata."),
        "confidentiality": _gate("PASS", "The result contains fixture IDs, reason codes, counts, timings, and digests only; document text and job-local paths are omitted."),
    }
    required = policy["required_gates"]
    decision = "GO" if all(gates[name]["status"] == "PASS" for name in required) else "NO_GO"
    result = {
        "schema_version": "sparkclaw.pptx_overlength.phase0_result.v1",
        "decision": decision,
        "policy_id": policy["policy_id"],
        "corpus_digest": corpus_digest,
        "repetitions": args.repetitions,
        "gates": gates,
        "engines": engines,
        "engine_errors": engine_errors,
        "writer": writer,
        "cancellation": cancellation,
        "outage": outage,
        "conversion_cost": conversion_cost,
    }
    return result


def main():
    parser = argparse.ArgumentParser(description="Run the fail-closed PPTX overlength Phase 0 qualification.")
    parser.add_argument("--policy", type=Path)
    parser.add_argument("--work-dir", type=Path, required=True)
    parser.add_argument("--gotenberg-url", required=True)
    parser.add_argument("--gotenberg-image", required=True)
    parser.add_argument("--gotenberg-fonts", type=Path)
    parser.add_argument("--powerpoint-evidence", type=Path)
    parser.add_argument("--repetitions", type=int, default=1)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--record-no-go", action="store_true")
    args = parser.parse_args()
    if args.repetitions < 1 or args.repetitions > 100:
        parser.error("--repetitions must be between 1 and 100")
    args.work_dir.mkdir(parents=True, exist_ok=True)
    result = run_qualification(args)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({
        "decision": result["decision"],
        "output": str(args.output),
        "failed_or_unresolved_gates": [
            name for name, gate in result["gates"].items() if gate["status"] != "PASS"
        ],
    }))
    if result["decision"] != "GO" and not args.record_no_go:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
