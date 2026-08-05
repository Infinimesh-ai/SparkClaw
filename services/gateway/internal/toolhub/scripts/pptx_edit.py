
import json, sys, os
try:
    from pptx import Presentation
except Exception:
    print(json.dumps({"error":"PPTX adapter requires python-pptx"}))
    sys.exit(0)

req = json.load(sys.stdin)
replacements = req.get("replacements") or []
counts = {item["Find"] if "Find" in item else item.get("find"): 0 for item in replacements}

def find_value(item):
    return item["Find"] if "Find" in item else item.get("find", "")

def replace_value(item):
    return item["Replace"] if "Replace" in item else item.get("replace", "")

def distribute_replacement(text, weights):
    total = sum(weights)
    if total <= 0:
        return [text] + [""] * (len(weights) - 1)
    boundaries = []
    cumulative = 0
    for weight in weights[:-1]:
        cumulative += weight
        boundaries.append(round(len(text) * cumulative / total))
    parts = []
    left = 0
    for index in range(len(weights)):
        right = boundaries[index] if index < len(boundaries) else len(text)
        parts.append(text[left:right])
        left = right
    return parts

def replace_span_in_paragraph(paragraph, start, end, replacement):
    runs = list(paragraph.runs)
    if not runs:
        return False
    full = "".join(run.text for run in runs)
    if start < 0 or end <= start or end > len(full):
        return False
    offsets = []
    cursor = 0
    for run in runs:
        next_cursor = cursor + len(run.text)
        offsets.append((cursor, next_cursor))
        cursor = next_cursor
    start_index = next(index for index, (_, right) in enumerate(offsets) if start < right)
    end_index = next(index for index, (_, right) in enumerate(offsets) if end <= right)
    start_left, _ = offsets[start_index]
    end_left, _ = offsets[end_index]
    prefix = runs[start_index].text[:start - start_left]
    suffix = runs[end_index].text[end - end_left:]
    affected = offsets[start_index:end_index + 1]
    weights = [max(0, min(end, right) - max(start, left)) for left, right in affected]
    parts = distribute_replacement(replacement, weights)
    for offset, index in enumerate(range(start_index, end_index + 1)):
        runs[index].text = (prefix if index == start_index else "") + parts[offset] + (suffix if index == end_index else "")
    return True

def replace_in_text_frame(tf):
    total = 0
    for paragraph in tf.paragraphs:
        for item in replacements:
            find = find_value(item)
            repl = replace_value(item)
            if not find:
                continue
            full = "".join(run.text for run in paragraph.runs)
            starts = []
            cursor = 0
            while True:
                start = full.find(find, cursor)
                if start < 0:
                    break
                starts.append(start)
                cursor = start + len(find)
            for start in reversed(starts):
                if not replace_span_in_paragraph(paragraph, start, start + len(find), repl):
                    raise ValueError("PPTX replacement span became stale")
                counts[find] = counts.get(find, 0) + 1
                total += 1
    return total

try:
    prs = Presentation(req["path"])
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                total += replace_in_text_frame(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        total += replace_in_text_frame(cell.text_frame)
    missing = [find for find, count in counts.items() if count == 0]
    if missing:
        print(json.dumps({"error":"find text was not matched: " + ", ".join(repr(x) for x in missing)}))
        sys.exit(0)
    os.makedirs(os.path.dirname(req["output_path"]), exist_ok=True)
    prs.save(req["output_path"])
    print(json.dumps({"replacements":total,"slides":len(prs.slides),"bytes":os.path.getsize(req["output_path"]),"details":[{"find":k,"count":v} for k,v in counts.items()]}))
except Exception as e:
    print(json.dumps({"error":str(e)}))
