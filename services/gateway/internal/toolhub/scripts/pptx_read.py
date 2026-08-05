import base64, json, re, sys
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    print(json.dumps({"error":"PPTX reader requires python-pptx"}))
    sys.exit(0)

req = json.load(sys.stdin)
max_bytes = int(req.get("max_bytes") or 20000)

def trim(text):
    return " ".join(str(text or "").split())

def enum_value(value):
    return str(value) if value is not None else ""

def shape_alt_text(shape):
    try:
        nodes = shape._element.xpath(".//p:cNvPr")
        if nodes:
            return str(nodes[0].get("descr") or nodes[0].get("title") or "")
    except Exception:
        pass
    return ""

def shape_geometry(shape, z_order, parent_group=""):
    placeholder_type = ""
    try:
        if shape.is_placeholder:
            placeholder_type = enum_value(shape.placeholder_format.type)
    except Exception:
        pass
    return {
        "shape_type": enum_value(shape.shape_type),
        "placeholder_type": placeholder_type,
        "x": int(getattr(shape, "left", 0) or 0),
        "y": int(getattr(shape, "top", 0) or 0),
        "width": int(getattr(shape, "width", 0) or 0),
        "height": int(getattr(shape, "height", 0) or 0),
        "z_order": int(z_order),
        "rotation": float(getattr(shape, "rotation", 0) or 0),
        "alternative_text": shape_alt_text(shape),
        "parent_group": parent_group,
    }

def color_value(color_format):
    try:
        if color_format.type is None:
            return ""
        if color_format.rgb is not None:
            return "#" + str(color_format.rgb)
    except Exception:
        pass
    try:
        if color_format.theme_color is not None:
            return "theme:" + enum_value(color_format.theme_color)
    except Exception:
        pass
    return ""

def shape_fill(shape):
    try:
        fill = shape.fill
        fill_type = enum_value(fill.type)
        color = color_value(fill.fore_color) if fill.type is not None else ""
        return {"type": fill_type, "color": color, "transparency": int(getattr(fill, "transparency", 0) or 0)}
    except Exception:
        return {"type": "", "color": "", "transparency": 0}

def shape_line(shape):
    try:
        return {"color": color_value(shape.line.color), "width": int(getattr(shape.line, "width", 0) or 0)}
    except Exception:
        return {"color": "", "width": 0}

def length_pt(value):
    try:
        return round(float(value.pt), 2) if value is not None else None
    except Exception:
        return None

def paragraph_bullet(paragraph):
    try:
        properties = paragraph._p.pPr
        if properties is None:
            return {"state": "inherited"}
        for child in properties:
            tag = str(child.tag)
            if tag.endswith("}buNone"):
                return {"state": "none"}
            if tag.endswith("}buChar"):
                return {"state": "character", "character": str(child.get("char") or "")}
            if tag.endswith("}buAutoNum"):
                return {"state": "auto_number", "scheme": str(child.get("type") or "")}
        return {"state": "inherited"}
    except Exception:
        return {"state": "unknown"}

def run_hyperlink(run):
    try:
        return str(run.hyperlink.address or "")
    except Exception:
        return ""

def run_style(run):
    font = run.font
    properties = run._r.get_or_add_rPr()
    return {
        "font_name": str(font.name or ""),
        "font_size_pt": length_pt(font.size),
        "bold": font.bold,
        "italic": font.italic,
        "underline": enum_value(font.underline),
        "font_color": color_value(font.color),
        "language": str(properties.get("lang") or ""),
        "alternative_language": str(properties.get("altLang") or ""),
        "baseline": str(properties.get("baseline") or ""),
        "hyperlink": run_hyperlink(run),
    }

def shape_text_structure(shape, slide_index, shape_index, path):
    if not getattr(shape, "has_text_frame", False):
        return {"editable": False, "paragraphs": [], "unsupported": ["no_text_frame"]}
    paragraphs = []
    unsupported = []
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
        runs = []
        for run_index, run in enumerate(paragraph.runs, start=1):
            runs.append({
                "path": "%s.paragraph[%d].run[%d]" % (path, paragraph_index, run_index),
                "index": run_index,
                "text": str(run.text or ""),
                "style": run_style(run),
            })
        fields = paragraph._p.xpath("./a:fld")
        if fields:
            unsupported.append("field")
        paragraphs.append({
            "path": "%s.paragraph[%d]" % (path, paragraph_index),
            "index": paragraph_index,
            "text": str(paragraph.text or ""),
            "level": int(paragraph.level or 0),
            "alignment": enum_value(paragraph.alignment),
            "space_before_pt": length_pt(paragraph.space_before),
            "space_after_pt": length_pt(paragraph.space_after),
            "line_spacing": str(paragraph.line_spacing or ""),
            "bullet": paragraph_bullet(paragraph),
            "soft_breaks": len(paragraph._p.xpath("./a:br")),
            "runs": runs,
        })
    return {
        "path": path + ".text_frame",
        "slide_index": slide_index,
        "shape_index": shape_index,
        "editable": len(unsupported) == 0,
        "paragraphs": paragraphs,
        "unsupported": sorted(set(unsupported)),
    }

def first_text_run(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if trim(run.text):
                return run
    return None

def shape_effective_font_size(shape):
    if not getattr(shape, "has_text_frame", False):
        return 18.0
    sizes = [
        float(run.font.size.pt)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if trim(run.text) and run.font.size is not None
    ]
    return max(sizes) if sizes else 18.0

def visual_text_units(text):
    units = 0.0
    for char in str(text or ""):
        code = ord(char)
        if char.isspace():
            units += 0.35
        elif code >= 0x2E80:
            units += 1.0
        elif char.isupper():
            units += 0.7
        elif char.islower() or char.isdigit():
            units += 0.55
        else:
            units += 0.5
    return round(units, 2)

def shape_text_style(shape):
    if not getattr(shape, "has_text_frame", False):
        return {}
    text_frame = shape.text_frame
    run = first_text_run(shape)
    font = run.font if run is not None else None
    size = None
    if font is not None and font.size is not None:
        size = round(float(font.size.pt), 2)
    name = str(font.name or "") if font is not None else ""
    color = color_value(font.color) if font is not None else ""
    bold = font.bold if font is not None else None
    text = trim(shape.text)
    usable_width = max(0, int(shape.width) - int(text_frame.margin_left or 0) - int(text_frame.margin_right or 0))
    capacity = None
    fit_ratio = None
    units = visual_text_units(text)
    effective_size = shape_effective_font_size(shape)
    if usable_width > 0:
        capacity = round((usable_width / 12700.0) / effective_size * 0.94, 2)
        fit_ratio = round(units / capacity, 3) if capacity > 0 else None
    return {
        "font_name": name,
        "font_size_pt": size,
        "bold": bold,
        "font_color": color,
        "word_wrap": bool(text_frame.word_wrap) if text_frame.word_wrap is not None else None,
        "auto_size": enum_value(text_frame.auto_size),
        "margin_left": int(text_frame.margin_left or 0),
        "margin_right": int(text_frame.margin_right or 0),
        "margin_top": int(text_frame.margin_top or 0),
        "margin_bottom": int(text_frame.margin_bottom or 0),
        "visual_units": units,
        "single_line_capacity_visual_units": capacity,
        "single_line_fit_ratio": fit_ratio,
    }

def layout_shape_record(slide_index, shape_index, path, shape, geometry, group_child_index=0):
    text = trim(shape.text) if getattr(shape, "has_text_frame", False) else ""
    structure = shape_text_structure(shape, slide_index, shape_index, path)
    record = dict(geometry)
    record.update({
        "slide_index": slide_index,
        "shape_index": shape_index,
        "path": path,
        "name": str(getattr(shape, "name", "") or ""),
        "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
        "text": text,
        "fill": shape_fill(shape),
        "line": shape_line(shape),
        "text_style": shape_text_style(shape),
        "editable": bool(text) and bool(structure.get("editable")) and not group_child_index,
    })
    if group_child_index:
        record["group_child_index"] = group_child_index
        record["editable"] = False
    return record

def vertical_center(record):
    return int(record["y"]) + int(record["height"]) / 2.0

def horizontal_center(record):
    return int(record["x"]) + int(record["width"]) / 2.0

def vertical_overlap(left, right):
    top = max(int(left["y"]), int(right["y"]))
    bottom = min(int(left["y"]) + int(left["height"]), int(right["y"]) + int(right["height"]))
    return max(0, bottom - top)

def derive_companion_groups(slide_index, records):
    groups = []
    top_level = [record for record in records if not record.get("parent_group")]
    backgrounds = [
        record for record in top_level
        if not record.get("text") and record.get("fill", {}).get("color") and int(record.get("width") or 0) > 0
    ]
    texts = [record for record in top_level if record.get("text")]
    used = set()
    for background in backgrounds:
        bg_left = int(background["x"])
        bg_right = bg_left + int(background["width"])
        bg_center = vertical_center(background)
        row_texts = [
            record for record in texts
            if record["shape_index"] not in used
            and int(record["y"]) <= bg_center <= int(record["y"]) + int(record["height"])
            and vertical_overlap(background, record) >= min(int(background["height"]), int(record["height"])) * 0.45
        ]
        labels = [
            record for record in row_texts
            if int(record["x"]) >= bg_left - int(background["width"]) * 0.03
            and int(record["x"]) + int(record["width"]) <= bg_right + int(background["width"]) * 0.03
        ]
        if not labels:
            continue
        label = min(labels, key=lambda record: horizontal_center(record))
        bodies = [
            record for record in row_texts
            if record["shape_index"] != label["shape_index"]
            and int(record["x"]) > int(label["x"])
            and int(record["x"]) <= bg_right + int(background["width"]) * 0.03
            and int(record["x"]) + int(record["width"]) > bg_right + int(background["width"]) * 0.20
        ]
        if not bodies:
            continue
        body = min(bodies, key=lambda record: abs(vertical_center(record) - bg_center))
        group_id = "slide:%d:band:%d" % (slide_index, int(background["shape_index"]))
        group = {
            "id": group_id,
            "kind": "label_body_band",
            "confidence": "high",
            "slide_index": slide_index,
            "background_shape_index": int(background["shape_index"]),
            "label_shape_index": int(label["shape_index"]),
            "body_shape_index": int(body["shape_index"]),
            "path": "presentation.slide[%d].companion_group[%d]" % (slide_index, len(groups) + 1),
        }
        groups.append(group)
        for record, role in ((background, "background"), (label, "label"), (body, "body")):
            record["companion_group_id"] = group_id
            record["companion_role"] = role
            used.add(record["shape_index"])
    return groups

page_marker_pattern = re.compile(r"(?<!\d)(\d+)\s*/\s*(\d+)(?!\d)")

extracted_bytes = 0

def append_line(lines, text):
    global extracted_bytes
    extracted_bytes += len(text.encode("utf-8")) + (1 if lines else 0)
    if extracted_bytes > max_bytes:
        print(json.dumps({"content":"", "truncated":True, "extracted_bytes":extracted_bytes}))
        sys.exit(0)
    lines.append(text)

try:
    prs = Presentation(req["path"])
    slides = []
    lines = []
    images = []
    resources = []
    notes = []
    hyperlinks = []
    charts = []
    slide_layouts = []
    layout_inventory = []
    layout_shapes = []
    companion_groups = []
    page_markers = []
    layout_warnings = []
    resource_keys = set()

    representative_slides = {}
    for slide_index, slide in enumerate(prs.slides, start=1):
        try:
            part_name = str(slide.slide_layout.part.partname)
        except Exception:
            part_name = ""
        representative_slides.setdefault(part_name, []).append(slide_index)
    for layout in prs.slide_layouts:
        try:
            part_name = str(layout.part.partname)
        except Exception:
            part_name = ""
        placeholders = []
        try:
            for placeholder in layout.placeholders:
                placeholders.append({
                    "name": str(getattr(placeholder, "name", "") or ""),
                    "placeholder_index": int(placeholder.placeholder_format.idx),
                    "role": enum_value(placeholder.placeholder_format.type),
                })
        except Exception:
            placeholders = []
        layout_inventory.append({
            "layout_ref": "layout:" + part_name,
            "name": str(getattr(layout, "name", "") or ""),
            "part_name": part_name,
            "placeholder_roles": placeholders,
            "representative_slide_refs": ["slide:%d" % value for value in representative_slides.get(part_name, [])[:3]],
        })

    def register_picture(slide, slide_index, shape, shape_index, path, geometry):
        try:
            blob = bytes(shape.image.blob)
            rel_id = str(shape._pic.blipFill.blip.rEmbed or "")
            related = slide.part.related_part(rel_id) if rel_id else None
            part_name = str(getattr(related, "partname", ""))
            content_type = str(getattr(related, "content_type", "application/octet-stream"))
        except Exception:
            return
        resource_key = "slide:%d:%s:%s" % (slide_index, rel_id, part_name)
        images.append({
            "kind": "image",
            "resource_key": resource_key,
            "parent_path": "presentation.slide[%d]" % slide_index,
            "location": dict(geometry, slide_index=slide_index, shape_index=shape_index, path=path),
            "source": {"parser": "python_pptx", "relationship_id": rel_id, "part_name": part_name},
            "content_type": content_type,
        })
        if resource_key not in resource_keys:
            resource_keys.add(resource_key)
            resources.append({
                "key": resource_key,
                "kind": "image",
                "content_type": content_type,
                "data_base64": base64.b64encode(blob).decode("ascii"),
            })

    def collect_hyperlinks(shape, slide_index, shape_index, path):
        if not getattr(shape, "has_text_frame", False):
            return
        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
            for run_index, run in enumerate(paragraph.runs, start=1):
                try:
                    target = run.hyperlink.address
                except Exception:
                    target = None
                if target:
                    hyperlinks.append({
                        "kind": "hyperlink", "text": trim(run.text), "target": str(target),
                        "location": {"slide_index": slide_index, "shape_index": shape_index, "paragraph_index": paragraph_index, "run_index": run_index, "path": path},
                        "source": {"parser": "python_pptx"},
                    })

    def chart_record(slide_index, shape_index, path, shape, geometry):
        chart = shape.chart
        title = ""
        try:
            if chart.has_title:
                title = trim(chart.chart_title.text_frame.text)
        except Exception:
            pass
        series = []
        try:
            for item in chart.series:
                series.append({"name": str(item.name or ""), "values": [value for value in item.values]})
        except Exception:
            pass
        categories = []
        try:
            if chart.plots and chart.plots[0].categories:
                categories = [str(value) for value in chart.plots[0].categories]
        except Exception:
            pass
        return {
            "kind": "chart", "title": title, "chart_type": enum_value(chart.chart_type), "categories": categories, "series": series,
            "location": dict(geometry, slide_index=slide_index, shape_index=shape_index, path=path),
            "source": {"parser": "python_pptx"},
        }

    for slide_index, slide in enumerate(prs.slides, start=1):
        items = []
        slide_shape_records = []
        try:
            layout_name = str(getattr(slide.slide_layout, "name", "") or "")
            layout_part = str(getattr(slide.slide_layout.part, "partname", "") or "")
        except Exception:
            layout_name = ""
            layout_part = ""
        layout_ref = "layout:" + layout_part
        slide_layouts.append({
            "slide_index": slide_index,
            "layout_ref": layout_ref,
            "name": layout_name,
            "part_name": layout_part,
            "path": "presentation.slide[%d]" % slide_index,
        })
        for shape_index, shape in enumerate(slide.shapes, start=1):
            path = "presentation.slide[%d].shape[%d]" % (slide_index, shape_index)
            geometry = shape_geometry(shape, shape_index)
            shape_record = layout_shape_record(slide_index, shape_index, path, shape, geometry)
            slide_shape_records.append(shape_record)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                register_picture(slide, slide_index, shape, shape_index, path, geometry)
            collect_hyperlinks(shape, slide_index, shape_index, path)
            if getattr(shape, "has_chart", False):
                chart = chart_record(slide_index, shape_index, path, shape, geometry)
                charts.append(chart)
                items.append(dict(chart, type="chart", shape_index=shape_index, path=path))
            if getattr(shape, "has_text_frame", False) and trim(shape.text):
                text_structure = shape_text_structure(shape, slide_index, shape_index, path)
                marker_match = page_marker_pattern.search(trim(shape.text))
                if marker_match and int(shape.top) >= int(prs.slide_height * 0.75):
                    page_markers.append({
                        "slide_index": slide_index,
                        "shape_index": shape_index,
                        "current": int(marker_match.group(1)),
                        "declared_total": int(marker_match.group(2)),
                        "actual_total": len(prs.slides),
                        "path": path,
                    })
                items.append({
                    "shape_index": shape_index,
                    "type": "text",
                    "text": trim(shape.text),
                    "path": path,
                    "layout": geometry,
                    "editable": bool(text_structure.get("editable")),
                    "text_structure": text_structure,
                })
            if getattr(shape, "has_table", False):
                rows = []
                for row_index, row in enumerate(shape.table.rows, start=1):
                    rows.append({"index": row_index, "cells": [trim(cell.text) for cell in row.cells]})
                items.append({"shape_index": shape_index, "type": "table", "rows": rows, "path": path, "layout": geometry})
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_index, child in enumerate(shape.shapes, start=1):
                    child_path = path + ".group_shape[%d]" % child_index
                    child_geometry = shape_geometry(child, child_index, path)
                    slide_shape_records.append(layout_shape_record(slide_index, shape_index, child_path, child, child_geometry, child_index))
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        register_picture(slide, slide_index, child, shape_index, child_path, child_geometry)
                    collect_hyperlinks(child, slide_index, shape_index, child_path)
                    if getattr(child, "has_text_frame", False) and trim(child.text):
                        child_structure = shape_text_structure(child, slide_index, shape_index, child_path)
                        items.append({
                            "shape_index": shape_index, "group_child_index": child_index, "parent_group": path,
                            "type": "text", "text": trim(child.text), "path": child_path, "layout": child_geometry, "editable": False,
                            "text_structure": child_structure,
                        })
        slide_groups = derive_companion_groups(slide_index, slide_shape_records)
        layout_shapes.extend(slide_shape_records)
        companion_groups.extend(slide_groups)
        slides.append({
            "index": slide_index,
            "template_ref": "slide:%d" % slide_index,
            "layout_ref": layout_ref,
            "layout_name": layout_name,
            "layout_part": layout_part,
            "items": items,
        })
        if items:
            append_line(lines, "Slide %d:" % slide_index)
            for item in items:
                if item["type"] == "text":
                    append_line(lines, item["text"])
                elif item["type"] == "table":
                    for row in item["rows"]:
                        append_line(lines, "\t".join(row["cells"]))
        try:
            notes_text = trim(slide.notes_slide.notes_text_frame.text)
            if notes_text:
                notes.append({
                    "kind": "speaker_notes", "text": notes_text,
                    "location": {"slide_index": slide_index, "path": "presentation.slide[%d].notes" % slide_index},
                    "source": {"parser": "python_pptx", "part_name": str(slide.notes_slide.part.partname)},
                })
        except Exception:
            pass

    for marker in page_markers:
        if marker["current"] != marker["slide_index"] or marker["declared_total"] != marker["actual_total"]:
            layout_warnings.append(
                "slide %d page marker %d/%d does not match physical position %d/%d"
                % (marker["slide_index"], marker["current"], marker["declared_total"], marker["slide_index"], marker["actual_total"])
            )

    content = "\n".join(lines).strip()
    raw = content.encode("utf-8")
    truncated = len(raw) > max_bytes
    if truncated:
        content = raw[:max_bytes].decode("utf-8", errors="ignore")
    print(json.dumps({
        "content": content,
        "truncated": truncated,
        "extracted_bytes": extracted_bytes,
        "resources": resources,
        "document": {
            "schema_version": "document_read_v1",
            "format": "pptx",
            "source": "python_pptx",
            "slides": slides,
            "enrichment": {
                "schema_version": "document_enrichment_v1",
                "assets": {"images": images, "charts": charts, "embedded_objects": []},
                "annotations": {"comments": [], "notes": notes, "hyperlinks": hyperlinks},
                "layout": {
                    "sections": [],
                    "page_settings": [{"part": "presentation", "width": int(prs.slide_width), "height": int(prs.slide_height)}],
                    "slide_layouts": slide_layouts,
                    "layout_inventory": layout_inventory,
                    "merged_ranges": [],
                    "shapes": layout_shapes,
                    "companion_groups": companion_groups,
                    "page_markers": page_markers,
                },
                "extensions": {"status": "deferred", "parts": []},
                "coverage": {"content": "complete", "assets": "partial", "annotations": "partial", "layout": "partial", "extensions": "deferred"},
                "warnings": layout_warnings,
            },
            "stats": {
                "slides": len(slides), "images": len(images), "charts": len(charts), "notes": len(notes), "hyperlinks": len(hyperlinks),
                "layout_shapes": len(layout_shapes), "slide_layouts": len(layout_inventory), "companion_groups": len(companion_groups), "page_markers": len(page_markers),
            }
        }
    }, ensure_ascii=False))
except Exception as e:
    print(json.dumps({"error":str(e)}))
