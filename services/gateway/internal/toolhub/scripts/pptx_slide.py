
import copy
import json
import math
import os
import re
import sys
try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except Exception:
    print(json.dumps({"error":"PPTX slide adapter requires python-pptx"}))
    sys.exit(0)

req = json.load(sys.stdin)
op = req.get("operation")
EMU_PER_PT = 12700.0
WRAP_WIDTH_FACTOR = 0.90
LINE_HEIGHT_FACTOR = 1.20
MAX_STANDALONE_WRAP_LINES = 4

def positive_index(value, name):
    idx = int(value or 0)
    if idx <= 0:
        raise ValueError("%s must be a positive 1-based integer" % name)
    return idx

def slide_at(prs, idx):
    if idx < 1 or idx > len(prs.slides):
        raise ValueError("slide_index out of range: %s" % idx)
    return prs.slides[idx - 1]

def delete_slide(prs, idx):
    slide = slide_at(prs, idx)
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[idx - 1]
    rel_id = slide_id.rId
    slide_id_list.remove(slide_id)
    prs.part.drop_rel(rel_id)

def fill_text_placeholders(slide, title, body):
    title = str(title or "")
    body = str(body or "")
    if title and slide.shapes.title is not None:
        slide.shapes.title.text = title
    if body:
        for shape in slide.placeholders:
            if shape == slide.shapes.title:
                continue
            if hasattr(shape, "text_frame"):
                shape.text = body
                return
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                shape.text = body
                return

def normalized_text(value):
    return " ".join(str(value or "").split())

def visual_text_units(text):
    units = 0.0
    for char in text:
        code = ord(char)
        if char.isspace():
            units += 0.35
        elif code >= 0x2E80:
            units += 1.0
        elif char.isupper():
            units += 0.7
        elif char.islower() or char.isdigit():
            units += 0.55
        else:
            units += 0.5
    return max(units, 1.0)

def fitted_single_line_size(shape, text, max_size_pt):
    text_frame = shape.text_frame
    usable_width_pt = max(
        1.0,
        (shape.width - text_frame.margin_left - text_frame.margin_right) / 12700.0,
    )
    return min(max_size_pt, usable_width_pt / visual_text_units(text) * 0.94)

def shape_font_size(shape):
    text_frame = shape.text_frame
    runs = [run for paragraph in text_frame.paragraphs for run in paragraph.runs if run.text]
    explicit_sizes = [run.font.size.pt for run in runs if run.font.size is not None]
    return max(explicit_sizes) if explicit_sizes else 18.0

def text_fits_single_line(shape, size_pt=None):
    if size_pt is None:
        size_pt = shape_font_size(shape)
    return fitted_single_line_size(shape, shape.text_frame.text, size_pt) >= size_pt - 0.25

def logical_text_lines(text):
    return re.split(r"[\r\n\v]", str(text or ""))

def text_line_capacity(shape, size_pt=None):
    if size_pt is None:
        size_pt = shape_font_size(shape)
    text_frame = shape.text_frame
    usable_width_pt = max(
        1.0,
        (shape.width - text_frame.margin_left - text_frame.margin_right) / EMU_PER_PT,
    )
    return max(1.0, usable_width_pt / max(size_pt, 1.0) * WRAP_WIDTH_FACTOR)

def wrapped_line_count(shape, text=None, size_pt=None):
    if text is None:
        text = shape.text_frame.text
    capacity = text_line_capacity(shape, size_pt)
    return max(1, sum(max(1, int(math.ceil(visual_text_units(line) / capacity))) for line in logical_text_lines(text)))

def required_text_height(shape, text=None, size_pt=None):
    if size_pt is None:
        size_pt = shape_font_size(shape)
    text_frame = shape.text_frame
    margins_pt = (text_frame.margin_top + text_frame.margin_bottom) / EMU_PER_PT
    lines = wrapped_line_count(shape, text, size_pt)
    return int(math.ceil((margins_pt + lines * size_pt * LINE_HEIGHT_FACTOR) * EMU_PER_PT))

def text_fits_wrapped(shape, size_pt=None):
    return required_text_height(shape, size_pt=size_pt) <= int(shape.height)

def measured_text_flow(shape, size_pt=None):
    explicit_breaks = len(logical_text_lines(shape.text_frame.text)) > 1
    if not explicit_breaks and text_fits_single_line(shape, size_pt):
        return "single"
    if text_fits_wrapped(shape, size_pt):
        return "wrapped"
    return "overflow"

def apply_measured_text_flow(shape, size_pt=None):
    flow = measured_text_flow(shape, size_pt)
    if flow == "single":
        return flow
    shape.text_frame.word_wrap = True
    return flow

def text_fits_current_flow(shape, size_pt=None):
    if size_pt is None:
        size_pt = shape_font_size(shape)
    lines = logical_text_lines(shape.text_frame.text)
    if len(lines) == 1 and text_fits_single_line(shape, size_pt):
        return True
    if shape.text_frame.word_wrap is False:
        if any(fitted_single_line_size(shape, line, size_pt) < size_pt - 0.25 for line in lines):
            return False
        return required_text_height(shape, size_pt=size_pt) <= int(shape.height)
    return text_fits_wrapped(shape, size_pt)

def shape_uses_multiple_lines(shape, size_pt=None):
    return len(logical_text_lines(shape.text_frame.text)) > 1 or wrapped_line_count(shape, size_pt=size_pt) > 1

def shape_bounds(shape):
    word_wrap = None
    if getattr(shape, "has_text_frame", False):
        word_wrap = shape.text_frame.word_wrap
    return {
        "x": int(shape.left),
        "y": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
        "font_size_pt": round(shape_font_size(shape), 2) if getattr(shape, "has_text_frame", False) and normalized_text(shape.text) else None,
        "word_wrap": bool(word_wrap) if word_wrap is not None else None,
    }

def vertical_overlap(left, right):
    top = max(int(left.top), int(right.top))
    bottom = min(int(left.top + left.height), int(right.top + right.height))
    return max(0, bottom - top)

def horizontal_overlap(left, right):
    start = max(int(left.left), int(right.left))
    end = min(int(left.left + left.width), int(right.left + right.width))
    return max(0, end - start)

def shape_has_fill(shape):
    try:
        return shape.fill.type is not None
    except Exception:
        return False

def derive_band_groups(slide):
    groups = []
    used = set()
    backgrounds = [
        (index, shape) for index, shape in enumerate(slide.shapes, start=1)
        if shape_has_fill(shape) and not normalized_text(getattr(shape, "text", "")) and int(shape.width) > 0
    ]
    texts = [
        (index, shape) for index, shape in enumerate(slide.shapes, start=1)
        if getattr(shape, "has_text_frame", False) and normalized_text(shape.text)
    ]
    for background_index, background in backgrounds:
        bg_left = int(background.left)
        bg_right = int(background.left + background.width)
        bg_center = int(background.top + background.height / 2)
        row_texts = [
            (index, shape) for index, shape in texts
            if index not in used
            and int(shape.top) <= bg_center <= int(shape.top + shape.height)
            and vertical_overlap(background, shape) >= min(int(background.height), int(shape.height)) * 0.45
        ]
        labels = [
            (index, shape) for index, shape in row_texts
            if int(shape.left) >= bg_left - int(background.width) * 0.03
            and int(shape.left + shape.width) <= bg_right + int(background.width) * 0.03
        ]
        if not labels:
            continue
        label_index, label = min(labels, key=lambda item: int(item[1].left + item[1].width / 2))
        bodies = [
            (index, shape) for index, shape in row_texts
            if index != label_index
            and int(shape.left) > int(label.left)
            and int(shape.left) <= bg_right + int(background.width) * 0.03
            and int(shape.left + shape.width) > bg_right + int(background.width) * 0.20
        ]
        if not bodies:
            continue
        body_index, body = min(bodies, key=lambda item: abs(int(item[1].top + item[1].height / 2) - bg_center))
        groups.append({
            "background_index": background_index,
            "background": background,
            "label_index": label_index,
            "label": label,
            "body_index": body_index,
            "body": body,
        })
        used.update((background_index, label_index, body_index))
    return groups

def shape_is_contained(background, shape, tolerance):
    return (
        int(shape.left) >= int(background.left) - tolerance
        and int(shape.top) >= int(background.top) - tolerance
        and int(shape.left + shape.width) <= int(background.left + background.width) + tolerance
        and int(shape.top + shape.height) <= int(background.top + background.height) + tolerance
    )

def derive_card_groups(slide):
    groups = []
    backgrounds = [
        (index, shape) for index, shape in enumerate(slide.shapes, start=1)
        if shape_has_fill(shape)
        and not normalized_text(getattr(shape, "text", ""))
        and int(shape.width) >= 914400
        and int(shape.height) >= 457200
    ]
    texts = [
        (index, shape) for index, shape in enumerate(slide.shapes, start=1)
        if getattr(shape, "has_text_frame", False) and normalized_text(shape.text)
    ]
    for background_index, background in backgrounds:
        tolerance = max(12700, int(min(background.width, background.height) * 0.04))
        children = [(index, shape) for index, shape in texts if shape_is_contained(background, shape, tolerance)]
        if len(children) < 2:
            continue
        title_index, title = min(children, key=lambda item: (int(item[1].top), int(item[1].left)))
        body_index, body = max(children, key=lambda item: (int(item[1].top), int(item[1].height)))
        if title_index == body_index or int(body.top) <= int(title.top):
            continue
        companions = []
        for index, shape in enumerate(slide.shapes, start=1):
            if index == background_index or normalized_text(getattr(shape, "text", "")) or not shape_has_fill(shape):
                continue
            if shape_is_contained(background, shape, tolerance):
                companions.append((index, shape))
        groups.append({
            "background_index": background_index,
            "background": background,
            "title_index": title_index,
            "title": title,
            "body_index": body_index,
            "body": body,
            "text_indexes": {index for index, _ in children},
            "companions": companions,
        })
    return groups

def peer_band_family(groups, selected_indexes):
    selected = [group for group in groups if group["body_index"] in selected_indexes]
    if not selected:
        return []
    anchor = selected[0]
    tolerance = max(12700, int(anchor["background"].width * 0.04))
    family = [
        group for group in groups
        if abs(int(group["background"].left) - int(anchor["background"].left)) <= tolerance
        and abs(int(group["background"].width) - int(anchor["background"].width)) <= tolerance
        and abs(int(group["body"].left) - int(anchor["body"].left)) <= tolerance
    ]
    return family if len(family) >= 2 else []

def peer_card_family(groups, selected_indexes):
    selected = [group for group in groups if group["body_index"] in selected_indexes]
    if not selected:
        return []
    anchor = selected[0]
    background = anchor["background"]
    body = anchor["body"]
    tolerance = max(12700, int(min(background.width, background.height) * 0.08))
    body_left_offset = int(body.left - background.left)
    body_top_offset = int(body.top - background.top)
    family = [
        group for group in groups
        if abs(int(group["background"].top) - int(background.top)) <= tolerance
        and abs(int(group["background"].width) - int(background.width)) <= tolerance
        and abs(int(group["background"].height) - int(background.height)) <= tolerance
        and abs(int(group["body"].left - group["background"].left) - body_left_offset) <= tolerance
        and abs(int(group["body"].top - group["background"].top) - body_top_offset) <= tolerance
    ]
    ordered = sorted(family, key=lambda group: int(group["background"].left))
    for current, following in zip(ordered, ordered[1:]):
        if int(current["background"].left + current["background"].width) > int(following["background"].left) + tolerance:
            return []
    return family if len(family) >= 2 else []

def safe_right_boundary(prs, slide, shape, excluded_indexes):
    boundary = int(prs.slide_width - min(prs.slide_width * 0.05, 457200))
    current_right = int(shape.left + shape.width)
    for index, other in enumerate(slide.shapes, start=1):
        if index in excluded_indexes or other is shape or vertical_overlap(shape, other) <= 0:
            continue
        other_left = int(other.left)
        if other_left >= current_right and other_left < boundary:
            boundary = other_left - 91440
    return max(current_right, boundary)

def safe_bottom_boundary(prs, slide, shape, excluded_indexes):
    boundary = int(prs.slide_height - min(prs.slide_height * 0.05, 342900))
    current_bottom = int(shape.top + shape.height)
    for index, other in enumerate(slide.shapes, start=1):
        if index in excluded_indexes or other is shape or horizontal_overlap(shape, other) <= 0:
            continue
        other_top = int(other.top)
        if other_top >= current_bottom and other_top < boundary:
            boundary = other_top - 91440
    return max(current_bottom, boundary)

def geometry_change(index, before, after):
    return {"shape_index": index, "before": before, "after": after}

def apply_coordinated_band_layout(prs, slide, groups, selected_indexes):
    family = peer_band_family(groups, selected_indexes)
    if not family:
        return [], set()
    excluded = set()
    for group in family:
        excluded.update((group["background_index"], group["label_index"], group["body_index"]))
    common_right = min(safe_right_boundary(prs, slide, group["body"], excluded) for group in family)
    for group in family:
        background = group["background"]
        body = group["body"]
        new_background_width = int(body.left - background.left)
        new_body_width = int(common_right - body.left)
        if new_background_width <= 0 or new_body_width <= 0:
            raise ValueError("coordinated slide layout could not establish non-overlapping label and body columns")
        background.width = new_background_width
        body.width = new_body_width
    target_body_height = max(
        max(int(group["body"].height) for group in family),
        max(required_text_height(group["body"]) for group in family),
    )
    bottom_padding = max(
        int(group["background"].top + group["background"].height - group["body"].top - group["body"].height)
        for group in family
    )
    target_background_height = max(
        max(int(group["background"].height) for group in family),
        max(int(group["body"].top - group["background"].top) + target_body_height + bottom_padding for group in family),
    )
    ordered = sorted(family, key=lambda group: int(group["background"].top))
    for current, following in zip(ordered, ordered[1:]):
        if int(current["background"].top) + target_background_height > int(following["background"].top) - 91440:
            raise ValueError("updated text is too tall for the coordinated peer-row layout")
    if int(ordered[-1]["background"].top) + target_background_height > safe_bottom_boundary(prs, slide, ordered[-1]["background"], excluded):
        raise ValueError("updated text is too tall for the coordinated peer-row layout")
    wrapped = set()
    for group in family:
        background = group["background"]
        label = group["label"]
        body = group["body"]
        background.height = target_background_height
        body.height = target_body_height
        if int(label.height) < target_body_height:
            label.height = target_body_height
        flow = apply_measured_text_flow(body)
        if flow == "overflow":
            raise ValueError("updated text does not fit the coordinated peer-row layout")
        if shape_uses_multiple_lines(body):
            wrapped.add(group["body_index"])
    return family, wrapped

def apply_coordinated_card_layout(prs, slide, groups, selected_indexes):
    family = peer_card_family(groups, selected_indexes)
    if not family:
        return [], set()
    members = set()
    for group in family:
        members.add(group["background_index"])
        members.update(group["text_indexes"])
        members.update(index for index, _ in group["companions"])
    target_body_height = max(
        max(int(group["body"].height) for group in family),
        max(required_text_height(group["body"]) for group in family),
    )
    bottom_padding = max(
        int(group["background"].top + group["background"].height - group["body"].top - group["body"].height)
        for group in family
    )
    target_background_height = max(
        max(int(group["background"].height) for group in family),
        max(int(group["body"].top - group["background"].top) + target_body_height + bottom_padding for group in family),
    )
    for group in family:
        if int(group["background"].top) + target_background_height > safe_bottom_boundary(prs, slide, group["background"], members):
            raise ValueError("updated text is too tall for the coordinated card layout")
    wrapped = set()
    for group in family:
        background = group["background"]
        body = group["body"]
        original_background_top = int(background.top)
        original_background_bottom = int(background.top + background.height)
        background.height = target_background_height
        body.height = target_body_height
        flow = apply_measured_text_flow(body)
        if flow == "overflow":
            raise ValueError("updated text does not fit the coordinated card layout")
        if shape_uses_multiple_lines(body):
            wrapped.add(group["body_index"])
        tolerance = max(12700, int(min(background.width, background.height) * 0.04))
        for _, companion in group["companions"]:
            top_aligned = abs(int(companion.top) - original_background_top) <= tolerance
            bottom_aligned = abs(int(companion.top + companion.height) - original_background_bottom) <= tolerance
            if top_aligned and bottom_aligned:
                companion.height = int(background.top + background.height - companion.top)
    return family, wrapped

def fit_shape_without_collision(prs, slide, shape_index, excluded_indexes):
    shape = slide.shapes[shape_index - 1]
    original_height = int(shape.height)
    explicit_breaks = len(logical_text_lines(shape.text_frame.text)) > 1
    flow = measured_text_flow(shape)
    if flow != "overflow":
        if flow == "wrapped":
            shape.text_frame.word_wrap = True
        return flow
    if not explicit_breaks:
        right = safe_right_boundary(prs, slide, shape, excluded_indexes)
        if right > int(shape.left + shape.width):
            shape.width = int(right - shape.left)
        flow = measured_text_flow(shape)
        if flow == "single":
            shape.text_frame.word_wrap = False
            return flow
        if flow == "wrapped":
            shape.text_frame.word_wrap = True
            return flow
    line_count = wrapped_line_count(shape)
    size_pt = shape_font_size(shape)
    text_frame = shape.text_frame
    original_usable_height_pt = max(
        0.0,
        (original_height - text_frame.margin_top - text_frame.margin_bottom) / EMU_PER_PT,
    )
    original_line_capacity = int(original_usable_height_pt // max(size_pt * LINE_HEIGHT_FACTOR, 1.0))
    if line_count > max(MAX_STANDALONE_WRAP_LINES, original_line_capacity):
        raise ValueError("updated text is too long for its slide shape after multi-line layout; shorten the text")
    required_height = required_text_height(shape)
    bottom = safe_bottom_boundary(prs, slide, shape, excluded_indexes)
    if int(shape.top) + required_height > bottom:
        raise ValueError("updated text is too tall for its slide shape without overlapping nearby content")
    shape.height = required_height
    shape.text_frame.word_wrap = True
    if not text_fits_current_flow(shape):
        raise ValueError("updated text does not fit its slide shape after multi-line layout")
    return "wrapped"

def validate_slide_layout(prs, slide, updated_indexes, band_groups, card_groups, before):
    for index, shape in enumerate(slide.shapes, start=1):
        if int(shape.left) < 0 or int(shape.top) < 0 or int(shape.left + shape.width) > int(prs.slide_width) or int(shape.top + shape.height) > int(prs.slide_height):
            raise ValueError("slide shape %s is outside the presentation canvas" % index)
    for index in updated_indexes:
        shape = slide.shapes[index - 1]
        if not text_fits_current_flow(shape):
            raise ValueError("updated text does not fit slide shape %s" % index)
    for group in band_groups:
        background = group["background"]
        body = group["body"]
        if int(background.left + background.width) > int(body.left):
            raise ValueError("coordinated label background overlaps body shape %s" % group["body_index"])
        if int(body.top + body.height) > int(background.top + background.height) + 12700:
            raise ValueError("coordinated peer-row body extends below its background shape %s" % group["background_index"])
    for group in card_groups:
        background = group["background"]
        body = group["body"]
        if not shape_is_contained(background, body, 12700):
            raise ValueError("coordinated card body extends beyond background shape %s" % group["background_index"])
    peer_font_uniform = True
    for family in (band_groups, card_groups):
        if not family:
            continue
        body_heights = {int(group["body"].height) for group in family}
        body_fonts = {round(shape_font_size(group["body"]), 2) for group in family}
        if len(body_heights) != 1:
            raise ValueError("coordinated peer text boxes are not geometrically uniform")
        if len(body_fonts) != 1:
            peer_font_uniform = False
        for group in family:
            body_index = group["body_index"]
            if round(shape_font_size(group["body"]), 2) != before[body_index]["font_size_pt"]:
                raise ValueError("coordinated peer body font size changed")
    return {
        "updated_text_fits": True,
        "wrapped_text_fits": True,
        "canvas_bounds": True,
        "companion_non_overlap": True,
        "peer_font_uniform": peer_font_uniform,
        "peer_font_preserved": True,
        "peer_geometry_uniform": True,
    }

def page_marker_warnings(prs):
    warnings = []
    marker_pattern = re.compile(r"(?<!\d)(\d+)\s*/\s*(\d+)(?!\d)")
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if int(shape.top) < int(prs.slide_height * 0.75):
                continue
            value = normalized_text(shape.text)
            match = marker_pattern.search(value)
            if not match:
                continue
            current, declared_total = int(match.group(1)), int(match.group(2))
            if current != slide_index or declared_total != len(prs.slides):
                warnings.append(
                    "slide %d page marker %d/%d does not match physical position %d/%d"
                    % (slide_index, current, declared_total, slide_index, len(prs.slides))
                )
    return warnings

def paragraph_has_fields(paragraph):
    return bool(paragraph._p.xpath("./a:fld"))

def copy_run_properties(source, destination):
    source_properties = copy.deepcopy(source._r.get_or_add_rPr())
    current_properties = destination._r.get_or_add_rPr()
    destination._r.remove(current_properties)
    destination._r.insert(0, source_properties)

def distribute_text_across_runs(paragraph, text):
    runs = list(paragraph.runs)
    if not runs:
        runs = [paragraph.add_run()]
    original_lengths = [len(run.text) for run in runs]
    original_total = sum(original_lengths)
    if original_total <= 0:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
        return
    boundaries = []
    cumulative = 0
    for length in original_lengths[:-1]:
        cumulative += length
        boundaries.append(round(len(text) * cumulative / original_total))
    start = 0
    for index, run in enumerate(runs):
        end = boundaries[index] if index < len(boundaries) else len(text)
        run.text = text[start:end]
        start = end

def distribute_replacement(text, weights):
    total = sum(weights)
    if total <= 0:
        return [text] + [""] * (len(weights) - 1)
    boundaries = []
    cumulative = 0
    for weight in weights[:-1]:
        cumulative += weight
        boundaries.append(round(len(text) * cumulative / total))
    parts = []
    left = 0
    for index in range(len(weights)):
        right = boundaries[index] if index < len(boundaries) else len(text)
        parts.append(text[left:right])
        left = right
    return parts

def replace_exact_span(shape, find, replacement):
    if not find:
        raise ValueError("exact_span update requires find")
    matches = []
    match_count = 0
    for paragraph in shape.text_frame.paragraphs:
        if paragraph_has_fields(paragraph):
            raise ValueError("PPTX text fields are not editable without loss")
        full = "".join(run.text for run in paragraph.runs)
        count = full.count(find)
        if count > 0:
            match_count += count
            matches.append((paragraph, full.find(find)))
    if match_count != 1:
        raise ValueError("exact_span find must match exactly once within one paragraph")
    paragraph, start = matches[0]
    runs = list(paragraph.runs)
    if not runs:
        raise ValueError("exact_span target has no editable runs")
    end = start + len(find)
    offsets = []
    cursor = 0
    for run in runs:
        next_cursor = cursor + len(run.text)
        offsets.append((cursor, next_cursor))
        cursor = next_cursor
    start_index = next(index for index, (_, right) in enumerate(offsets) if start < right)
    end_index = next(index for index, (_, right) in enumerate(offsets) if end <= right)
    start_left, _ = offsets[start_index]
    end_left, _ = offsets[end_index]
    prefix = runs[start_index].text[:start - start_left]
    suffix = runs[end_index].text[end - end_left:]
    affected = offsets[start_index:end_index + 1]
    weights = [max(0, min(end, right) - max(start, left)) for left, right in affected]
    parts = distribute_replacement(replacement, weights)
    for offset, index in enumerate(range(start_index, end_index + 1)):
        runs[index].text = (prefix if index == start_index else "") + parts[offset] + (suffix if index == end_index else "")

def rewrite_shape_text(shape, text, break_mode):
    paragraphs = list(shape.text_frame.paragraphs)
    if any(paragraph_has_fields(paragraph) for paragraph in paragraphs):
        raise ValueError("PPTX text fields are not editable without loss")
    lines = text.split("\n")
    if len(paragraphs) > 1:
        if break_mode != "paragraph" or len(lines) != len(paragraphs):
            raise ValueError("multi-paragraph shapes require break_mode=paragraph and one line per existing paragraph")
        for paragraph, line in zip(paragraphs, lines):
            distribute_text_across_runs(paragraph, line)
        return
    paragraph = paragraphs[0]
    for child in list(paragraph._p):
        if child.tag.endswith("}br"):
            paragraph._p.remove(child)
    distribute_text_across_runs(paragraph, lines[0])
    if len(lines) <= 1:
        return
    if break_mode != "soft_break":
        raise ValueError("single-paragraph multiline text requires break_mode=soft_break")
    source_run = paragraph.runs[-1] if paragraph.runs else paragraph.add_run()
    for line in lines[1:]:
        paragraph.add_line_break()
        run = paragraph.add_run()
        copy_run_properties(source_run, run)
        run.text = line

def replace_text_preserving_style(shape, update):
    text = re.sub(r"\r\n?|\v", "\n", str(update.get("text") or ""))
    if not text.strip():
        raise ValueError("updated shape text must not be empty")
    mode = normalized_text(update.get("mode") or "rewrite_shape").lower()
    if mode == "exact_span":
        replace_exact_span(shape, str(update.get("find") or ""), text)
        return
    if mode != "rewrite_shape":
        raise ValueError("PPTX text update mode must be rewrite_shape or exact_span")
    break_mode = normalized_text(update.get("break_mode") or ("paragraph" if len(shape.text_frame.paragraphs) > 1 else "soft_break")).lower()
    rewrite_shape_text(shape, text, break_mode)

def update_slide(prs, slide, updates, layout_policy):
    if not isinstance(updates, list) or not updates:
        raise ValueError("updates must be a non-empty array")
    layout_policy = normalized_text(layout_policy or "coordinated").lower()
    if layout_policy not in ("preserve", "coordinated"):
        raise ValueError("layout_policy must be preserve or coordinated")
    seen = set()
    before = {index: shape_bounds(shape) for index, shape in enumerate(slide.shapes, start=1)}
    for update in updates:
        if not isinstance(update, dict):
            raise ValueError("each slide update must be an object")
        shape_index = positive_index(update.get("shape_index"), "shape_index")
        if shape_index in seen:
            raise ValueError("shape_index is duplicated: %s" % shape_index)
        seen.add(shape_index)
        if shape_index > len(slide.shapes):
            raise ValueError("shape_index out of range: %s" % shape_index)
        shape = slide.shapes[shape_index - 1]
        if not getattr(shape, "has_text_frame", False):
            raise ValueError("shape_index does not identify a text shape: %s" % shape_index)
        expected = normalized_text(update.get("old_text"))
        current = normalized_text(shape.text)
        if not expected or current != expected:
            raise ValueError("old_text does not match slide shape %s" % shape_index)
        replace_text_preserving_style(shape, update)

    band_groups = derive_band_groups(slide)
    card_groups = derive_card_groups(slide)
    coordinated_bands = []
    coordinated_cards = []
    wrapped = set()
    if layout_policy == "coordinated":
        coordinated_bands, band_wrapped = apply_coordinated_band_layout(prs, slide, band_groups, seen)
        coordinated_cards, card_wrapped = apply_coordinated_card_layout(prs, slide, card_groups, seen)
        wrapped.update(band_wrapped)
        wrapped.update(card_wrapped)
        family_body_indexes = {
            group["body_index"] for group in coordinated_bands + coordinated_cards
        }
        excluded = set()
        for group in coordinated_bands:
            excluded.update((group["background_index"], group["label_index"], group["body_index"]))
        for group in coordinated_cards:
            excluded.add(group["background_index"])
            excluded.update(group["text_indexes"])
            excluded.update(index for index, _ in group["companions"])
        for shape_index in seen - family_body_indexes:
            flow = fit_shape_without_collision(prs, slide, shape_index, excluded)
            if flow == "wrapped":
                wrapped.add(shape_index)
    else:
        for shape_index in seen:
            shape = slide.shapes[shape_index - 1]
            if not text_fits_current_flow(shape):
                raise ValueError("updated text does not fit preserve layout_policy; use coordinated or shorten the text")
            if shape_uses_multiple_lines(shape):
                wrapped.add(shape_index)

    checks = validate_slide_layout(prs, slide, seen, coordinated_bands, coordinated_cards, before)
    changes = []
    for shape_index, shape in enumerate(slide.shapes, start=1):
        after = shape_bounds(slide.shapes[shape_index - 1])
        if before[shape_index] != after:
            changes.append(geometry_change(shape_index, before[shape_index], after))
    adjusted = {change["shape_index"] for change in changes}
    return {
        "updated_shapes": len(seen),
        "fitted_shapes": 0,
        "wrapped_shapes": len(wrapped),
        "wrapped_shape_indexes": sorted(wrapped),
        "layout_policy": layout_policy,
        "layout_adjusted_shapes": len(adjusted),
        "layout_adjusted_shape_indexes": sorted(adjusted),
        "layout_changes": changes,
        "layout_checks": checks,
        "companion_groups_used": len(coordinated_bands) + len(coordinated_cards),
    }

def move_last_slide_after(prs, after_slide_index):
    if after_slide_index < 0 or after_slide_index >= len(prs.slides):
        raise ValueError("after_slide_index out of range: %s" % after_slide_index)
    slide_id_list = prs.slides._sldIdLst
    new_slide_id = slide_id_list[-1]
    slide_id_list.remove(new_slide_id)
    slide_id_list.insert(after_slide_index, new_slide_id)

def duplicate_slide(prs, idx, after_slide_index=None):
    source = slide_at(prs, idx)
    if getattr(source, "has_notes_slide", False) and normalized_text(source.notes_slide.notes_text_frame.text):
        raise ValueError("template or duplicate slide contains speaker notes, which cannot be cloned without loss")
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in list(dest.shapes):
        element = shape.element
        element.getparent().remove(element)
    relationship_ids = {}
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if rel.is_external:
            relationship_ids[rel.rId] = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            relationship_ids[rel.rId] = dest.part.rels.get_or_add(rel.reltype, rel._target)
    relationship_attributes = (qn("r:embed"), qn("r:link"), qn("r:id"))
    for shape in source.shapes:
        element = copy.deepcopy(shape.element)
        for child in element.iter():
            for attribute in relationship_attributes:
                old_id = child.get(attribute)
                if old_id in relationship_ids:
                    child.set(attribute, relationship_ids[old_id])
        dest.shapes._spTree.insert_element_before(element, 'p:extLst')
    if after_slide_index is None:
        after_slide_index = idx
    move_last_slide_after(prs, after_slide_index)
    return dest

def layout_for_ref(prs, layout_ref):
    for layout in prs.slide_layouts:
        if "layout:" + str(layout.part.partname) == layout_ref:
            return layout
    raise ValueError("layout_ref is stale or does not belong to the current presentation")

def slide_index_for_ref(prs, slide_ref):
    match = re.fullmatch(r"slide:([1-9][0-9]*)", str(slide_ref or ""))
    if not match:
        raise ValueError("template_slide_ref is invalid")
    index = int(match.group(1))
    slide_at(prs, index)
    return index

try:
    prs = Presentation(req["path"])
    result = {
        "status": "pptx_version_written",
        "operation": op,
        "path": req["path"],
        "output_path": req["output_path"]
    }
    if op == "add_slide":
        before_count = len(prs.slides)
        after_slide_index = int(req.get("after_slide_index") or before_count)
        if after_slide_index < 1 or after_slide_index > before_count:
            raise ValueError("after_slide_index out of range: %s" % after_slide_index)
        layout_ref = str(req.get("layout_ref") or "")
        template_slide_ref = str(req.get("template_slide_ref") or "")
        if bool(layout_ref) == bool(template_slide_ref):
            raise ValueError("exactly one of layout_ref or template_slide_ref is required")
        if template_slide_ref:
            template_index = slide_index_for_ref(prs, template_slide_ref)
            slide = duplicate_slide(prs, template_index, after_slide_index)
            updates = req.get("template_updates") or []
            if updates:
                result.update(update_slide(prs, slide, updates, req.get("layout_policy") or "coordinated"))
            result["template_slide_ref"] = template_slide_ref
        else:
            slide = prs.slides.add_slide(layout_for_ref(prs, layout_ref))
            move_last_slide_after(prs, after_slide_index)
            fill_text_placeholders(slide, req.get("title"), req.get("body"))
            result["layout_ref"] = layout_ref
        result["slide_index"] = after_slide_index + 1
        result["inserted_slide_index"] = after_slide_index + 1
        result["after_slide_index"] = after_slide_index
        result["title"] = str(req.get("title") or "")
        result["body"] = str(req.get("body") or "")
        result["warnings"] = page_marker_warnings(prs)
    elif op == "update_slide":
        idx = positive_index(req.get("slide_index"), "slide_index")
        slide = slide_at(prs, idx)
        result.update(update_slide(prs, slide, req.get("updates"), req.get("layout_policy")))
        result["warnings"] = page_marker_warnings(prs)
        result["slide_index"] = idx
    elif op == "update_deck":
        slide_updates = req.get("slide_updates") or []
        if not isinstance(slide_updates, list) or not slide_updates:
            raise ValueError("slide_updates must be a non-empty array")
        seen_slides = set()
        aggregate = {
            "updated_slides": 0, "updated_shapes": 0, "wrapped_shapes": 0,
            "layout_adjusted_shapes": 0, "companion_groups_used": 0,
            "slide_indexes": [], "wrapped_shape_indexes": [],
            "layout_adjusted_shape_indexes": [], "layout_adjusted_targets": [],
            "layout_changes": [], "layout_checks": {},
        }
        for slide_update in slide_updates:
            if not isinstance(slide_update, dict):
                raise ValueError("each slide_updates item must be an object")
            idx = positive_index(slide_update.get("slide_index"), "slide_index")
            if idx in seen_slides:
                raise ValueError("slide_index is duplicated: %s" % idx)
            seen_slides.add(idx)
            current = update_slide(prs, slide_at(prs, idx), slide_update.get("updates"), slide_update.get("layout_policy") or "coordinated")
            aggregate["updated_slides"] += 1
            aggregate["slide_indexes"].append(idx)
            for key in ("updated_shapes", "wrapped_shapes", "layout_adjusted_shapes", "companion_groups_used"):
                aggregate[key] += int(current.get(key) or 0)
            aggregate["wrapped_shape_indexes"].extend(current.get("wrapped_shape_indexes") or [])
            aggregate["layout_adjusted_shape_indexes"].extend(current.get("layout_adjusted_shape_indexes") or [])
            for shape_index in current.get("layout_adjusted_shape_indexes") or []:
                aggregate["layout_adjusted_targets"].append({"slide_index": idx, "shape_index": shape_index})
            for change in current.get("layout_changes") or []:
                aggregate["layout_changes"].append(dict(change, slide_index=idx))
            aggregate["layout_checks"]["slide:%d" % idx] = current.get("layout_checks") or {}
        aggregate["warnings"] = page_marker_warnings(prs)
        result.update(aggregate)
    elif op == "duplicate_slide":
        idx = positive_index(req.get("slide_index"), "slide_index")
        duplicate_slide(prs, idx)
        result["slide_index"] = idx
        result["inserted_slide_index"] = idx + 1
        result["warnings"] = page_marker_warnings(prs)
    elif op == "delete_slide":
        idx = positive_index(req.get("slide_index"), "slide_index")
        if len(prs.slides) <= 1:
            raise ValueError("cannot delete the only slide")
        delete_slide(prs, idx)
        result["slide_index"] = idx
        result["warnings"] = page_marker_warnings(prs)
    else:
        raise ValueError("unsupported pptx operation: %s" % op)
    os.makedirs(os.path.dirname(req["output_path"]), exist_ok=True)
    prs.save(req["output_path"])
    result["slides"] = len(prs.slides)
    result["bytes"] = os.path.getsize(req["output_path"])
    print(json.dumps(result, ensure_ascii=False))
except Exception as e:
    print(json.dumps({"error":str(e)}, ensure_ascii=False))
