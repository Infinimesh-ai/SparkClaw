import json
import os
import sys

try:
    from pptx import Presentation
except Exception:
    print(json.dumps({"error":"PPTX slide adapter requires python-pptx"}))
    sys.exit(0)

from .clone import duplicate_slide, layout_for_ref, move_last_slide_after, slide_index_for_ref
from .errors import PPTXLayoutFitError
from .layout import page_marker_warnings
from .slides import delete_slide, fill_text_placeholders, positive_index, slide_at
from .update import update_slide

req = json.load(sys.stdin)
op = req.get("operation")


try:
    prs = Presentation(req["path"])
    result = {
        "status": "pptx_version_written",
        "operation": op,
        "path": req["path"],
        "output_path": req["output_path"]
    }
    if op == "add_slide":
        before_count = len(prs.slides)
        after_slide_index = int(req.get("after_slide_index") or before_count)
        if after_slide_index < 1 or after_slide_index > before_count:
            raise ValueError("after_slide_index out of range: %s" % after_slide_index)
        layout_ref = str(req.get("layout_ref") or "")
        template_slide_ref = str(req.get("template_slide_ref") or "")
        if bool(layout_ref) == bool(template_slide_ref):
            raise ValueError("exactly one of layout_ref or template_slide_ref is required")
        if template_slide_ref:
            template_index = slide_index_for_ref(prs, template_slide_ref)
            slide = duplicate_slide(prs, template_index, after_slide_index)
            updates = req.get("template_updates") or []
            if updates:
                result.update(update_slide(prs, slide, updates, req.get("layout_policy") or "coordinated"))
            result["template_slide_ref"] = template_slide_ref
        else:
            slide = prs.slides.add_slide(layout_for_ref(prs, layout_ref))
            move_last_slide_after(prs, after_slide_index)
            fill_text_placeholders(slide, req.get("title"), req.get("body"))
            result["layout_ref"] = layout_ref
        result["slide_index"] = after_slide_index + 1
        result["inserted_slide_index"] = after_slide_index + 1
        result["after_slide_index"] = after_slide_index
        result["title"] = str(req.get("title") or "")
        result["body"] = str(req.get("body") or "")
        result["warnings"] = page_marker_warnings(prs)
    elif op == "update_slide":
        idx = positive_index(req.get("slide_index"), "slide_index")
        slide = slide_at(prs, idx)
        result.update(update_slide(prs, slide, req.get("updates"), req.get("layout_policy")))
        result["warnings"] = page_marker_warnings(prs)
        result["slide_index"] = idx
    elif op == "update_deck":
        slide_updates = req.get("slide_updates") or []
        if not isinstance(slide_updates, list) or not slide_updates:
            raise ValueError("slide_updates must be a non-empty array")
        seen_slides = set()
        aggregate = {
            "updated_slides": 0, "updated_shapes": 0, "wrapped_shapes": 0,
            "layout_adjusted_shapes": 0, "companion_groups_used": 0,
            "slide_indexes": [], "wrapped_shape_indexes": [],
            "layout_adjusted_shape_indexes": [], "layout_adjusted_targets": [],
            "layout_changes": [], "layout_checks": {},
        }
        for slide_update in slide_updates:
            if not isinstance(slide_update, dict):
                raise ValueError("each slide_updates item must be an object")
            idx = positive_index(slide_update.get("slide_index"), "slide_index")
            if idx in seen_slides:
                raise ValueError("slide_index is duplicated: %s" % idx)
            seen_slides.add(idx)
            current = update_slide(prs, slide_at(prs, idx), slide_update.get("updates"), slide_update.get("layout_policy") or "coordinated")
            aggregate["updated_slides"] += 1
            aggregate["slide_indexes"].append(idx)
            for key in ("updated_shapes", "wrapped_shapes", "layout_adjusted_shapes", "companion_groups_used"):
                aggregate[key] += int(current.get(key) or 0)
            aggregate["wrapped_shape_indexes"].extend(current.get("wrapped_shape_indexes") or [])
            aggregate["layout_adjusted_shape_indexes"].extend(current.get("layout_adjusted_shape_indexes") or [])
            for shape_index in current.get("layout_adjusted_shape_indexes") or []:
                aggregate["layout_adjusted_targets"].append({"slide_index": idx, "shape_index": shape_index})
            for change in current.get("layout_changes") or []:
                aggregate["layout_changes"].append(dict(change, slide_index=idx))
            aggregate["layout_checks"]["slide:%d" % idx] = current.get("layout_checks") or {}
        aggregate["warnings"] = page_marker_warnings(prs)
        result.update(aggregate)
    elif op == "duplicate_slide":
        idx = positive_index(req.get("slide_index"), "slide_index")
        duplicate_slide(prs, idx)
        result["slide_index"] = idx
        result["inserted_slide_index"] = idx + 1
        result["warnings"] = page_marker_warnings(prs)
    elif op == "delete_slide":
        idx = positive_index(req.get("slide_index"), "slide_index")
        if len(prs.slides) <= 1:
            raise ValueError("cannot delete the only slide")
        delete_slide(prs, idx)
        result["slide_index"] = idx
        result["warnings"] = page_marker_warnings(prs)
    else:
        raise ValueError("unsupported pptx operation: %s" % op)
    os.makedirs(os.path.dirname(req["output_path"]), exist_ok=True)
    prs.save(req["output_path"])
    result["slides"] = len(prs.slides)
    result["bytes"] = os.path.getsize(req["output_path"])
    print(json.dumps(result, ensure_ascii=False))
except Exception as e:
    error = {"error":str(e)}
    if isinstance(e, PPTXLayoutFitError):
        error["error_code"] = "pptx_layout_fit_conflict"
    print(json.dumps(error, ensure_ascii=False))
