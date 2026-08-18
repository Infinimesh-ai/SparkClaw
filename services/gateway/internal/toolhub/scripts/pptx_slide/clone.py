import copy
import re

from pptx.oxml.ns import qn

from .slides import slide_at
from .text import normalized_text


def move_last_slide_after(prs, after_slide_index):
    if after_slide_index < 0 or after_slide_index >= len(prs.slides):
        raise ValueError("after_slide_index out of range: %s" % after_slide_index)
    slide_id_list = prs.slides._sldIdLst
    new_slide_id = slide_id_list[-1]
    slide_id_list.remove(new_slide_id)
    slide_id_list.insert(after_slide_index, new_slide_id)

def duplicate_slide(prs, idx, after_slide_index=None):
    source = slide_at(prs, idx)
    if getattr(source, "has_notes_slide", False) and normalized_text(source.notes_slide.notes_text_frame.text):
        raise ValueError("template or duplicate slide contains speaker notes, which cannot be cloned without loss")
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in list(dest.shapes):
        element = shape.element
        element.getparent().remove(element)
    relationship_ids = {}
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if rel.is_external:
            relationship_ids[rel.rId] = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            relationship_ids[rel.rId] = dest.part.rels.get_or_add(rel.reltype, rel._target)
    relationship_attributes = (qn("r:embed"), qn("r:link"), qn("r:id"))
    for shape in source.shapes:
        element = copy.deepcopy(shape.element)
        for child in element.iter():
            for attribute in relationship_attributes:
                old_id = child.get(attribute)
                if old_id in relationship_ids:
                    child.set(attribute, relationship_ids[old_id])
        dest.shapes._spTree.insert_element_before(element, 'p:extLst')
    if after_slide_index is None:
        after_slide_index = idx
    move_last_slide_after(prs, after_slide_index)
    return dest

def layout_for_ref(prs, layout_ref):
    for layout in prs.slide_layouts:
        if "layout:" + str(layout.part.partname) == layout_ref:
            return layout
    raise ValueError("layout_ref is stale or does not belong to the current presentation")

def slide_index_for_ref(prs, slide_ref):
    match = re.fullmatch(r"slide:([1-9][0-9]*)", str(slide_ref or ""))
    if not match:
        raise ValueError("template_slide_ref is invalid")
    index = int(match.group(1))
    slide_at(prs, index)
    return index
