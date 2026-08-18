import copy
import re

from pptx.oxml.ns import qn

from .text import normalized_text


def paragraph_has_fields(paragraph):
    return bool(paragraph._p.xpath("./a:fld"))

def copy_run_properties(source, destination):
    source_properties = copy.deepcopy(source._r.get_or_add_rPr())
    current_properties = destination._r.get_or_add_rPr()
    destination._r.remove(current_properties)
    destination._r.insert(0, source_properties)

def distribute_text_across_runs(paragraph, text):
    runs = list(paragraph.runs)
    if not runs:
        runs = [paragraph.add_run()]
    original_lengths = [len(run.text) for run in runs]
    original_total = sum(original_lengths)
    if original_total <= 0:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
        return
    boundaries = []
    cumulative = 0
    for length in original_lengths[:-1]:
        cumulative += length
        boundaries.append(round(len(text) * cumulative / original_total))
    start = 0
    for index, run in enumerate(runs):
        end = boundaries[index] if index < len(boundaries) else len(text)
        run.text = text[start:end]
        start = end

def distribute_replacement(text, weights):
    total = sum(weights)
    if total <= 0:
        return [text] + [""] * (len(weights) - 1)
    boundaries = []
    cumulative = 0
    for weight in weights[:-1]:
        cumulative += weight
        boundaries.append(round(len(text) * cumulative / total))
    parts = []
    left = 0
    for index in range(len(weights)):
        right = boundaries[index] if index < len(boundaries) else len(text)
        parts.append(text[left:right])
        left = right
    return parts

def replace_exact_span(shape, find, replacement):
    if not find:
        raise ValueError("exact_span update requires find")
    matches = []
    match_count = 0
    for paragraph in shape.text_frame.paragraphs:
        if paragraph_has_fields(paragraph):
            raise ValueError("PPTX text fields are not editable without loss")
        full = "".join(run.text for run in paragraph.runs)
        count = full.count(find)
        if count > 0:
            match_count += count
            matches.append((paragraph, full.find(find)))
    if match_count != 1:
        raise ValueError("exact_span find must match exactly once within one paragraph")
    paragraph, start = matches[0]
    runs = list(paragraph.runs)
    if not runs:
        raise ValueError("exact_span target has no editable runs")
    end = start + len(find)
    offsets = []
    cursor = 0
    for run in runs:
        next_cursor = cursor + len(run.text)
        offsets.append((cursor, next_cursor))
        cursor = next_cursor
    start_index = next(index for index, (_, right) in enumerate(offsets) if start < right)
    end_index = next(index for index, (_, right) in enumerate(offsets) if end <= right)
    start_left, _ = offsets[start_index]
    end_left, _ = offsets[end_index]
    prefix = runs[start_index].text[:start - start_left]
    suffix = runs[end_index].text[end - end_left:]
    affected = offsets[start_index:end_index + 1]
    weights = [max(0, min(end, right) - max(start, left)) for left, right in affected]
    parts = distribute_replacement(replacement, weights)
    for offset, index in enumerate(range(start_index, end_index + 1)):
        runs[index].text = (prefix if index == start_index else "") + parts[offset] + (suffix if index == end_index else "")

def rewrite_shape_text(shape, text, break_mode):
    paragraphs = list(shape.text_frame.paragraphs)
    if any(paragraph_has_fields(paragraph) for paragraph in paragraphs):
        raise ValueError("PPTX text fields are not editable without loss")
    lines = text.split("\n")
    if len(paragraphs) > 1:
        if break_mode != "paragraph" or len(lines) != len(paragraphs):
            raise ValueError("multi-paragraph shapes require break_mode=paragraph and one line per existing paragraph")
        for paragraph, line in zip(paragraphs, lines):
            distribute_text_across_runs(paragraph, line)
        return
    paragraph = paragraphs[0]
    for child in list(paragraph._p):
        if child.tag.endswith("}br"):
            paragraph._p.remove(child)
    distribute_text_across_runs(paragraph, lines[0])
    if len(lines) <= 1:
        return
    if break_mode != "soft_break":
        raise ValueError("single-paragraph multiline text requires break_mode=soft_break")
    source_run = paragraph.runs[-1] if paragraph.runs else paragraph.add_run()
    for line in lines[1:]:
        paragraph.add_line_break()
        run = paragraph.add_run()
        copy_run_properties(source_run, run)
        run.text = line

def replace_text_preserving_style(shape, update):
    text = re.sub(r"\r\n?|\v", "\n", str(update.get("text") or ""))
    if not text.strip():
        raise ValueError("updated shape text must not be empty")
    mode = normalized_text(update.get("mode") or "rewrite_shape").lower()
    if mode == "exact_span":
        replace_exact_span(shape, str(update.get("find") or ""), text)
        return
    if mode != "rewrite_shape":
        raise ValueError("PPTX text update mode must be rewrite_shape or exact_span")
    break_mode = normalized_text(update.get("break_mode") or ("paragraph" if len(shape.text_frame.paragraphs) > 1 else "soft_break")).lower()
    rewrite_shape_text(shape, text, break_mode)
