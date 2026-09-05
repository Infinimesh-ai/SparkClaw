import hashlib
import json
import math
import os
import re
import sys
import unicodedata
from io import BytesIO

try:
    import pypdfium2 as pdfium
    from PIL import Image
    from lxml import etree
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as exc:
    print(json.dumps({"error": "PPTX visual QA requires python-pptx, lxml, Pillow, and pypdfium2: %s" % exc, "error_code": "pptx_visual_qa_unavailable"}))
    sys.exit(0)


ANALYSIS_SCHEMA = "sparkclaw.pptx_render_analysis.v1"
DIAGNOSTIC_SCHEMA = "sparkclaw.pptx_diagnostic_facts.v1"
CONTEXT_SCHEMA = "sparkclaw.pptx_visual_repair_context.v1"
MAX_SHAPES_PER_PAGE = 128
MAX_FACTS_PER_PAGE = 128
MAX_TEXT_CHARS = 600


def fail(message, code="pptx_visual_qa_failed"):
    print(json.dumps({"error": str(message), "error_code": code}, ensure_ascii=False))
    sys.exit(0)


def normalize_text(value):
    return " ".join(unicodedata.normalize("NFKC", str(value or "")).split())


def bounded_text(value):
    value = str(value or "")
    return value[:MAX_TEXT_CHARS]


def enum_value(value):
    return str(value) if value is not None else ""


def color_value(color_format):
    try:
        if color_format.type is None:
            return ""
        if color_format.rgb is not None:
            return "#" + str(color_format.rgb)
    except Exception:
        pass
    try:
        if color_format.theme_color is not None:
            return "theme:" + enum_value(color_format.theme_color)
    except Exception:
        pass
    return ""


def length_pt(value):
    try:
        return round(float(value.pt), 2) if value is not None else None
    except Exception:
        return None


def region_milli(x, y, width, height, slide_width, slide_height):
    return [
        round(x * 1000 / slide_width),
        round(y * 1000 / slide_height),
        round(width * 1000 / slide_width),
        round(height * 1000 / slide_height),
    ]


def point_milli(point, slide_width, slide_height):
    return [round(point[0] * 1000 / slide_width), round(point[1] * 1000 / slide_height)]


def rotated_rectangle(x, y, width, height, rotation):
    points = [(x, y), (x + width, y), (x + width, y + height), (x, y + height)]
    if not rotation:
        return points
    radians = math.radians(rotation)
    cosine = math.cos(radians)
    sine = math.sin(radians)
    center_x = x + width / 2.0
    center_y = y + height / 2.0
    out = []
    for point_x, point_y in points:
        delta_x = point_x - center_x
        delta_y = point_y - center_y
        out.append((center_x + delta_x * cosine - delta_y * sine, center_y + delta_x * sine + delta_y * cosine))
    return out


def polygon_area(points):
    if len(points) < 3:
        return 0.0
    total = 0.0
    for index, point in enumerate(points):
        next_point = points[(index + 1) % len(points)]
        total += point[0] * next_point[1] - next_point[0] * point[1]
    return abs(total) / 2.0


def polygon_orientation(points):
    total = 0.0
    for index, point in enumerate(points):
        next_point = points[(index + 1) % len(points)]
        total += (next_point[0] - point[0]) * (next_point[1] + point[1])
    return -1.0 if total > 0 else 1.0


def polygon_intersection(subject, clip):
    output = list(subject)
    orientation = polygon_orientation(clip)

    def inside(point, edge_start, edge_end):
        cross = (edge_end[0] - edge_start[0]) * (point[1] - edge_start[1]) - (edge_end[1] - edge_start[1]) * (point[0] - edge_start[0])
        return cross * orientation >= -1e-7

    def intersection(start, end, edge_start, edge_end):
        delta_x = end[0] - start[0]
        delta_y = end[1] - start[1]
        edge_x = edge_end[0] - edge_start[0]
        edge_y = edge_end[1] - edge_start[1]
        denominator = delta_x * edge_y - delta_y * edge_x
        if abs(denominator) < 1e-12:
            return end
        factor = ((edge_start[0] - start[0]) * edge_y - (edge_start[1] - start[1]) * edge_x) / denominator
        return (start[0] + factor * delta_x, start[1] + factor * delta_y)

    for edge_index, edge_start in enumerate(clip):
        if not output:
            break
        edge_end = clip[(edge_index + 1) % len(clip)]
        input_points = output
        output = []
        previous = input_points[-1]
        for current in input_points:
            current_inside = inside(current, edge_start, edge_end)
            previous_inside = inside(previous, edge_start, edge_end)
            if current_inside:
                if not previous_inside:
                    output.append(intersection(previous, current, edge_start, edge_end))
                output.append(current)
            elif previous_inside:
                output.append(intersection(previous, current, edge_start, edge_end))
            previous = current
    return output


def shape_fill(shape):
    try:
        fill = shape.fill
        return {
            "type": enum_value(fill.type),
            "color": color_value(fill.fore_color) if fill.type is not None else "",
            "transparency": int(getattr(fill, "transparency", 0) or 0),
        }
    except Exception:
        return {"type": "", "color": "", "transparency": 0}


def shape_line(shape):
    try:
        return {"color": color_value(shape.line.color), "width": int(getattr(shape.line, "width", 0) or 0)}
    except Exception:
        return {"color": "", "width": 0}


def text_style(shape):
    if not getattr(shape, "has_text_frame", False):
        return {}
    first_run = None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if normalize_text(run.text):
                first_run = run
                break
        if first_run is not None:
            break
    font = first_run.font if first_run is not None else None
    return {
        "font_name": str(font.name or "") if font is not None else "",
        "font_size_pt": length_pt(font.size) if font is not None else None,
        "bold": font.bold if font is not None else None,
        "font_color": color_value(font.color) if font is not None else "",
        "word_wrap": bool(shape.text_frame.word_wrap) if shape.text_frame.word_wrap is not None else None,
        "auto_size": enum_value(shape.text_frame.auto_size),
        "margin_emu": [
            int(shape.text_frame.margin_left or 0),
            int(shape.text_frame.margin_top or 0),
            int(shape.text_frame.margin_right or 0),
            int(shape.text_frame.margin_bottom or 0),
        ],
    }


def group_child_transform(group):
    try:
        transform = group._element.grpSpPr.xfrm
        child_offset_x = int(transform.chOff.x)
        child_offset_y = int(transform.chOff.y)
        child_extent_x = int(transform.chExt.cx)
        child_extent_y = int(transform.chExt.cy)
        scale_x = float(group.width) / child_extent_x if child_extent_x else 1.0
        scale_y = float(group.height) / child_extent_y if child_extent_y else 1.0
        return child_offset_x, child_offset_y, scale_x, scale_y
    except Exception:
        return 0, 0, 1.0, 1.0


def collect_shapes(slide, slide_index, changed_indexes, changed_all, operation, slide_width, slide_height):
    records = []
    truncated = False

    def append_shape(shape, shape_ref, shape_index, z_order, parent_ref, x, y, width, height, rotation, editable):
        nonlocal truncated
        if len(records) >= MAX_SHAPES_PER_PAGE:
            truncated = True
            return False
        original_target_hash = shape_target_hash("", shape_ref, shape)
        full_text = normalize_text(shape.text) if getattr(shape, "has_text_frame", False) else ""
        polygon = rotated_rectangle(x, y, width, height, rotation)
        is_created = bool(operation == "add_slide" and changed_all and not parent_ref)
        records.append({
            "shape_ref": shape_ref,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "shape_type": enum_value(shape.shape_type),
            "role": enum_value(getattr(getattr(shape, "placeholder_format", None), "type", "")),
            "name": bounded_text(getattr(shape, "name", "")),
            "region_milli": region_milli(x, y, width, height, slide_width, slide_height),
            "polygon": polygon,
            "z_order": z_order,
            "rotation": round(float(rotation or 0), 3),
            "parent_group_ref": parent_ref,
            "text": bounded_text(full_text),
            "text_truncated": len(full_text) > MAX_TEXT_CHARS,
            "_full_text": full_text,
            "_target_hash": original_target_hash,
            "editable": bool(editable and full_text),
            "changed": bool(changed_all or shape_index in changed_indexes),
            "created": is_created,
            "edit_capabilities": [
                value for value, enabled in (
                    ("set_geometry", not parent_ref),
                    ("set_text_style", bool(editable and full_text)),
                    ("set_shape_style", not parent_ref),
                    ("place_above", not parent_ref),
                    ("place_below", not parent_ref),
                    ("rewrite_text", bool(editable and full_text)),
                    ("delete_generated_shape", is_created),
                ) if enabled
            ],
            "fill": shape_fill(shape),
            "line": shape_line(shape),
            "text_style": text_style(shape),
        })
        return True

    for shape_index, shape in enumerate(slide.shapes, start=1):
        shape_ref = "slide:%d:shape:%d" % (slide_index, shape_index)
        x = float(getattr(shape, "left", 0) or 0)
        y = float(getattr(shape, "top", 0) or 0)
        width = float(getattr(shape, "width", 0) or 0)
        height = float(getattr(shape, "height", 0) or 0)
        rotation = float(getattr(shape, "rotation", 0) or 0)
        if not append_shape(shape, shape_ref, shape_index, shape_index, "", x, y, width, height, rotation, True):
            break
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        child_offset_x, child_offset_y, scale_x, scale_y = group_child_transform(shape)
        for child_index, child in enumerate(shape.shapes, start=1):
            child_ref = shape_ref + ":child:%d" % child_index
            child_x = x + (float(getattr(child, "left", 0) or 0) - child_offset_x) * scale_x
            child_y = y + (float(getattr(child, "top", 0) or 0) - child_offset_y) * scale_y
            child_width = float(getattr(child, "width", 0) or 0) * scale_x
            child_height = float(getattr(child, "height", 0) or 0) * scale_y
            child_rotation = rotation + float(getattr(child, "rotation", 0) or 0)
            if not append_shape(child, child_ref, shape_index, shape_index * 1000 + child_index, shape_ref, child_x, child_y, child_width, child_height, child_rotation, False):
                break
        if truncated:
            break
    return records, truncated


def public_shape(record):
    return {key: value for key, value in record.items() if key not in ("polygon", "_full_text", "_target_hash")}


def shape_target_hash(candidate_sha256, shape_ref, shape):
    canonical = etree.tostring(shape._element, method="c14n", exclusive=True, with_comments=False)
    payload = shape_ref.encode("utf-8") + b"\0" + canonical
    return hashlib.sha256(payload).hexdigest()


def glyph_bounds_milli(text_page, pdf_width, pdf_height, shape_region):
    x, y, width, height = shape_region
    left = x * pdf_width / 1000.0
    right = (x + width) * pdf_width / 1000.0
    top = pdf_height - y * pdf_height / 1000.0
    bottom = pdf_height - (y + height) * pdf_height / 1000.0
    boxes = []
    for index in range(min(text_page.count_chars(), 10000)):
        try:
            char_left, char_bottom, char_right, char_top = text_page.get_charbox(index, loose=True)
        except Exception:
            continue
        center_x = (char_left + char_right) / 2.0
        center_y = (char_bottom + char_top) / 2.0
        if left <= center_x <= right and bottom <= center_y <= top:
            boxes.append((char_left, char_bottom, char_right, char_top))
    if not boxes:
        return []
    bound_left = min(box[0] for box in boxes)
    bound_bottom = min(box[1] for box in boxes)
    bound_right = max(box[2] for box in boxes)
    bound_top = max(box[3] for box in boxes)
    return [
        round(bound_left * 1000 / pdf_width),
        round((pdf_height - bound_top) * 1000 / pdf_height),
        round((bound_right - bound_left) * 1000 / pdf_width),
        round((bound_top - bound_bottom) * 1000 / pdf_height),
    ]


def clipping_fact(slide_index, record, text_page, page_text, pdf_width, pdf_height):
    expected = normalize_text(record.get("_full_text"))
    if not expected or not record.get("changed") or not record.get("editable"):
        return None
    x, y, width, height = record["region_milli"]
    left = x * pdf_width / 1000.0
    right = (x + width) * pdf_width / 1000.0
    top = pdf_height - y * pdf_height / 1000.0
    bottom = pdf_height - (y + height) * pdf_height / 1000.0
    try:
        observed = normalize_text(text_page.get_text_bounded(left, bottom, right, top))
    except Exception:
        observed = ""
    if expected == observed or expected in observed:
        return None
    status = "unavailable" if not page_text else "ambiguous"
    missing_start = 0
    while missing_start < min(len(expected), len(observed)) and expected[missing_start] == observed[missing_start]:
        missing_start += 1
    if not record.get("text_truncated"):
        if observed and expected.startswith(observed) and len(observed) >= max(3, math.ceil(len(expected) * 0.5)):
            status = "confirmed"
            missing_start = len(observed)
        elif observed:
            status = "observed"
    missing_spans = []
    if missing_start < len(expected):
        missing_spans.append({"start": missing_start, "end": len(expected), "text": bounded_text(expected[missing_start:])})
    return {
        "diagnostic_id": "diag-text-%d-%d" % (slide_index, record["shape_index"]),
        "kind": "text_clipping",
        "status": status,
        "shape_refs": [record["shape_ref"]],
        "evidence": {
            "text_frame_region_milli": record["region_milli"],
            "expected_text": bounded_text(expected),
            "observed_text": bounded_text(observed),
            "expected_text_truncated": bool(record.get("text_truncated")),
            "observed_text_truncated": len(observed) > MAX_TEXT_CHARS,
            "missing_spans": missing_spans,
            "rendered_glyph_bounds_milli": glyph_bounds_milli(text_page, pdf_width, pdf_height, record["region_milli"]),
            "coverage_status": "unavailable" if not page_text else "truncated" if record.get("text_truncated") else "complete",
        },
    }


def geometry_facts(slide_index, records, slide_width, slide_height, tolerance_milli):
    facts = []
    page_area = float(slide_width * slide_height)
    minimum_area = page_area * (float(tolerance_milli) / 1000.0) ** 2

    def append_fact(fact):
        if len(facts) >= MAX_FACTS_PER_PAGE:
            return False
        facts.append(fact)
        return True

    for record in records:
        if not record.get("changed"):
            continue
        polygon = record["polygon"]
        min_x = min(point[0] for point in polygon)
        min_y = min(point[1] for point in polygon)
        max_x = max(point[0] for point in polygon)
        max_y = max(point[1] for point in polygon)
        overflow = {
            "left": max(0.0, -min_x),
            "top": max(0.0, -min_y),
            "right": max(0.0, max_x - slide_width),
            "bottom": max(0.0, max_y - slide_height),
        }
        if any(value * 1000 / (slide_width if side in ("left", "right") else slide_height) > tolerance_milli for side, value in overflow.items()):
            if not append_fact({
                "diagnostic_id": "diag-canvas-%d-%d" % (slide_index, record["z_order"]),
                "kind": "off_canvas",
                "status": "confirmed",
                "shape_refs": [record["shape_ref"]],
                "evidence": {
                    "canvas_region_milli": [0, 0, 1000, 1000],
                    "shape_polygon_milli": [point_milli(point, slide_width, slide_height) for point in polygon],
                    "overflow_milli": {
                        "left": round(overflow["left"] * 1000 / slide_width),
                        "top": round(overflow["top"] * 1000 / slide_height),
                        "right": round(overflow["right"] * 1000 / slide_width),
                        "bottom": round(overflow["bottom"] * 1000 / slide_height),
                    },
                    "rotation": record["rotation"],
                    "parent_group_ref": record["parent_group_ref"],
                    "tolerance_milli": tolerance_milli,
                },
            }):
                return facts, True

    for left_index, left in enumerate(records):
        for right in records[left_index + 1:]:
            if not left.get("changed") and not right.get("changed"):
                continue
            if left["shape_ref"].startswith(right["shape_ref"] + ":") or right["shape_ref"].startswith(left["shape_ref"] + ":"):
                continue
            intersection = polygon_intersection(left["polygon"], right["polygon"])
            area = polygon_area(intersection)
            if area <= minimum_area:
                continue
            left_area = max(polygon_area(left["polygon"]), 1.0)
            right_area = max(polygon_area(right["polygon"]), 1.0)
            if not append_fact({
                "diagnostic_id": "diag-overlap-%d-%d-%d" % (slide_index, left["z_order"], right["z_order"]),
                "kind": "geometry_overlap",
                "status": "confirmed",
                "shape_refs": [left["shape_ref"], right["shape_ref"]],
                "evidence": {
                    "intersection_polygon_milli": [point_milli(point, slide_width, slide_height) for point in intersection],
                    "intersection_area_milli": round(area * 1000000 / page_area),
                    "overlap_ratio_milli": [round(area * 1000 / left_area), round(area * 1000 / right_area)],
                    "z_order": [left["z_order"], right["z_order"]],
                    "fill_transparency": [left["fill"]["transparency"], right["fill"]["transparency"]],
                    "group_ancestry": [left["parent_group_ref"], right["parent_group_ref"]],
                    "tolerance_milli": tolerance_milli,
                },
            }):
                return facts, True
    return facts, False


def raster_page(page, output_path, scale, max_page_pixels, max_png_bytes):
    bitmap = page.render(scale=scale, rotation=0)
    image = bitmap.to_pil().convert("RGB")
    pixel_count = image.width * image.height
    if pixel_count <= 0 or pixel_count > max_page_pixels:
        fail("rendered page pixel count exceeds the configured limit", "pptx_visual_qa_raster_limit")
    extrema = image.getextrema()
    uniform_black = all(low == 0 and high == 0 for low, high in extrema)
    uniform_white = all(low == 255 and high == 255 for low, high in extrema)
    if uniform_black:
        fail("rendered page is uniformly black", "pptx_visual_qa_invalid_png")
    buffer = BytesIO()
    image.save(buffer, format="PNG", compress_level=9, optimize=False)
    raw = buffer.getvalue()
    if len(raw) > max_png_bytes:
        fail("rendered PNG exceeds the configured byte limit", "pptx_visual_qa_raster_limit")
    with open(output_path, "wb") as handle:
        handle.write(raw)
    return {
        "png_path": output_path,
        "png_sha256": hashlib.sha256(raw).hexdigest(),
        "png_bytes": len(raw),
        "width": image.width,
        "height": image.height,
        "pixel_count": pixel_count,
        "uniform_white": uniform_white,
        "uniform_black": uniform_black,
    }


try:
    request = json.load(sys.stdin)
    candidate_path = os.path.abspath(request["path"])
    pdf_path = os.path.abspath(request["pdf_path"])
    output_dir = os.path.abspath(request["output_dir"])
    candidate_sha256 = str(request.get("candidate_sha256") or "")
    operation = str(request.get("operation") or "")
    selected = sorted(set(int(value) for value in request.get("slide_indexes") or []))
    changed_shape_indexes = {
        int(slide_index): set(int(value) for value in values)
        for slide_index, values in (request.get("changed_shape_indexes") or {}).items()
    }
    changed_all_slides = set(int(value) for value in request.get("changed_all_slides") or [])
    scale = float(request.get("raster_scale") or 1.5)
    max_pages = int(request.get("max_pages") or 100)
    max_changed_pages = int(request.get("max_changed_pages") or 20)
    max_page_pixels = int(request.get("max_page_pixels") or 20000000)
    max_png_bytes = int(request.get("max_png_bytes") or 12582912)
    tolerance_milli = int(request.get("diagnostic_tolerance_milli") or 2)

    if len(selected) > max_changed_pages:
        fail("changed slide count exceeds the configured limit", "pptx_visual_qa_page_selection")
    if any(slide_index not in selected for slide_index in changed_shape_indexes):
        fail("changed shape evidence references an unselected slide", "pptx_visual_qa_page_selection")
    if any(index <= 0 for indexes in changed_shape_indexes.values() for index in indexes):
        fail("changed shape evidence contains an invalid shape index", "pptx_visual_qa_page_selection")
    if any(slide_index not in selected for slide_index in changed_all_slides):
        fail("whole-slide change evidence references an unselected slide", "pptx_visual_qa_page_selection")
    os.makedirs(output_dir, exist_ok=True)

    presentation = Presentation(candidate_path)
    slide_count = len(presentation.slides)
    if slide_count <= 0 or slide_count > max_pages:
        fail("candidate slide count exceeds the configured limit", "pptx_visual_qa_page_count")
    if any(index <= 0 or index > slide_count for index in selected):
        fail("changed slide index is outside the candidate", "pptx_visual_qa_page_selection")

    pdf = pdfium.PdfDocument(pdf_path)
    if len(pdf) != slide_count:
        fail("rendered PDF page count does not match the PPTX slide count", "pptx_visual_qa_page_count")

    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    expected_aspect = slide_width / slide_height
    page_sizes = {}
    for slide_index in range(1, slide_count + 1):
        page = pdf[slide_index - 1]
        pdf_width, pdf_height = page.get_size()
        if not all(math.isfinite(value) and value > 0 for value in (pdf_width, pdf_height)):
            fail("rendered PDF page has invalid dimensions", "pptx_visual_qa_page_dimensions")
        if abs((pdf_width / pdf_height) - expected_aspect) / expected_aspect > 0.01:
            fail("rendered PDF page aspect ratio does not match the PPTX", "pptx_visual_qa_page_dimensions")
        page_sizes[slide_index] = (pdf_width, pdf_height)

    pages = []
    for slide_index in selected:
        page = pdf[slide_index - 1]
        pdf_width, pdf_height = page_sizes[slide_index]

        slide = presentation.slides[slide_index - 1]
        records, shapes_truncated = collect_shapes(
            slide,
            slide_index,
            changed_shape_indexes.get(slide_index, set()),
            slide_index in changed_all_slides,
            operation,
            slide_width,
            slide_height,
        )
        text_page = page.get_textpage()
        page_text = normalize_text(text_page.get_text_range())
        facts, facts_truncated = geometry_facts(slide_index, records, slide_width, slide_height, tolerance_milli)
        if not facts_truncated:
            for record in records:
                fact = clipping_fact(slide_index, record, text_page, page_text, pdf_width, pdf_height)
                if fact is None:
                    continue
                if len(facts) >= MAX_FACTS_PER_PAGE:
                    facts_truncated = True
                    break
                facts.append(fact)

        png_path = os.path.join(output_dir, "slide-%d.png" % slide_index)
        raster = raster_page(page, png_path, scale, max_page_pixels, max_png_bytes)
        pages.append({
            "slide_index": slide_index,
            "pdf_width": pdf_width,
            "pdf_height": pdf_height,
            "raster": raster,
            "structure": {
                "schema_version": CONTEXT_SCHEMA,
                "slide_index": slide_index,
                "page_region_milli": [0, 0, 1000, 1000],
                "shapes": [public_shape(record) for record in records],
                "truncated": shapes_truncated,
            },
            "targets": {
                record["shape_ref"]: record["_target_hash"]
                for record in records
                if not record["parent_group_ref"]
            },
            "diagnostics": {
                "schema_version": DIAGNOSTIC_SCHEMA,
                "candidate_sha256": candidate_sha256,
                "slide_index": slide_index,
                "coordinate_space": "region_milli",
                "facts": facts,
                "truncated": facts_truncated,
            },
        })

    print(json.dumps({
        "schema_version": ANALYSIS_SCHEMA,
        "candidate_sha256": candidate_sha256,
        "slide_count": slide_count,
        "slide_width": slide_width,
        "slide_height": slide_height,
        "pages": pages,
    }, ensure_ascii=False))
except Exception as exc:
    fail(exc)
