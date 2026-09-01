import hashlib
import json
import os
import re
import sys
import zipfile

try:
    from lxml import etree
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
except Exception as exc:
    print(json.dumps({"error": "PPTX visual repair requires python-pptx and lxml: %s" % exc, "error_code": "pptx_visual_repair_unavailable"}))
    sys.exit(0)


SHAPE_REF = re.compile(r"^slide:(\d+):shape:(\d+)$")
SHA256 = re.compile(r"^[0-9a-f]{64}$")
COLOR = re.compile(r"^#[0-9A-Fa-f]{6}$")
MAX_OPERATIONS = 8
MAX_REWRITE_CHARS = 1200


def fail(message, code="pptx_visual_repair_invalid"):
    print(json.dumps({"error": str(message), "error_code": code}, ensure_ascii=False))
    sys.exit(0)


def file_sha256(path):
    digest = hashlib.sha256()
    with open(path, "rb") as handle:
        for chunk in iter(lambda: handle.read(65536), b""):
            digest.update(chunk)
    return digest.hexdigest()


def target_hash(candidate_sha256, shape_ref, shape):
    canonical = etree.tostring(shape._element, method="c14n", exclusive=True, with_comments=False)
    payload = shape_ref.encode("utf-8") + b"\0" + canonical
    return hashlib.sha256(payload).hexdigest()


def canonical_part(name, raw):
    if name.endswith(".xml") or name.endswith(".rels"):
        parser = etree.XMLParser(resolve_entities=False, no_network=True, remove_blank_text=False)
        root = etree.fromstring(raw, parser=parser)
        return etree.tostring(root, method="c14n", exclusive=True, with_comments=False)
    return raw


def package_snapshot(path):
    snapshot = {}
    with zipfile.ZipFile(path, "r") as package:
        for name in package.namelist():
            if name.endswith("/"):
                continue
            raw = canonical_part(name, package.read(name))
            snapshot[name] = hashlib.sha256(raw).hexdigest()
    return snapshot


def parse_shape_ref(value, slide_index):
    match = SHAPE_REF.fullmatch(str(value or ""))
    if match is None or int(match.group(1)) != slide_index:
        fail("repair operation contains an invalid or cross-slide shape_ref")
    return int(match.group(2))


def require_shape(slide, shape_index):
    if shape_index <= 0 or shape_index > len(slide.shapes):
        fail("repair operation references a missing shape")
    return slide.shapes[shape_index - 1]


def require_int_list(value, length, minimum, maximum, field):
    if not isinstance(value, list) or len(value) != length or any(type(item) is not int for item in value):
        fail("%s must contain exactly %d integers" % (field, length))
    if any(item < minimum or item > maximum for item in value):
        fail("%s is outside the allowed range" % field)
    return value


def apply_geometry(shape, region, slide_width, slide_height):
    x, y, width, height = require_int_list(region, 4, 0, 1000, "region_milli")
    if width < 5 or height < 5 or x + width > 1000 or y + height > 1000:
        fail("region_milli must be a positive in-canvas rectangle")
    shape.left = round(x * slide_width / 1000)
    shape.top = round(y * slide_height / 1000)
    shape.width = round(width * slide_width / 1000)
    shape.height = round(height * slide_height / 1000)


def apply_text_style(shape, operation, slide_width, slide_height):
    if not getattr(shape, "has_text_frame", False):
        fail("set_text_style requires a text shape")
    allowed = {"op", "shape_ref", "font_size_pt", "alignment", "word_wrap", "margins_milli"}
    if any(key not in allowed for key in operation):
        fail("set_text_style contains an unsupported field")
    if "font_size_pt" in operation:
        size = operation["font_size_pt"]
        if not isinstance(size, (int, float)) or isinstance(size, bool) or size < 8 or size > 72:
            fail("font_size_pt must be between 8 and 72")
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(float(size))
    if "alignment" in operation:
        alignment = str(operation["alignment"])
        values = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
        if alignment not in values:
            fail("alignment is unsupported")
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = values[alignment]
    if "word_wrap" in operation:
        if type(operation["word_wrap"]) is not bool:
            fail("word_wrap must be boolean")
        shape.text_frame.word_wrap = operation["word_wrap"]
    if "margins_milli" in operation:
        left, top, right, bottom = require_int_list(operation["margins_milli"], 4, 0, 100, "margins_milli")
        shape.text_frame.margin_left = round(left * slide_width / 1000)
        shape.text_frame.margin_top = round(top * slide_height / 1000)
        shape.text_frame.margin_right = round(right * slide_width / 1000)
        shape.text_frame.margin_bottom = round(bottom * slide_height / 1000)


def apply_shape_style(shape, operation):
    allowed = {"op", "shape_ref", "fill_color", "line_color"}
    if any(key not in allowed for key in operation):
        fail("set_shape_style contains an unsupported field")
    if "fill_color" not in operation and "line_color" not in operation:
        fail("set_shape_style requires fill_color or line_color")
    if "fill_color" in operation:
        value = str(operation["fill_color"])
        if COLOR.fullmatch(value) is None:
            fail("fill_color must be #RRGGBB")
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(value[1:].upper())
    if "line_color" in operation:
        value = str(operation["line_color"])
        if COLOR.fullmatch(value) is None:
            fail("line_color must be #RRGGBB")
        shape.line.color.rgb = RGBColor.from_string(value[1:].upper())


def apply_rewrite(shape, operation):
    if not getattr(shape, "has_text_frame", False):
        fail("rewrite_text requires a text shape")
    text = operation.get("text")
    if not isinstance(text, str) or not text.strip() or len(text) > MAX_REWRITE_CHARS:
        fail("rewrite_text contains invalid text")
    first_run = None
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs:
            first_run = paragraph.runs[0]
            break
    if first_run is None:
        fail("rewrite_text requires an existing styled run")
    first_run.text = text
    seen = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run is first_run:
                seen = True
                continue
            if seen:
                run.text = ""


try:
    request = json.load(sys.stdin)
    input_path = os.path.abspath(request["path"])
    output_path = os.path.abspath(request["output_path"])
    expected_candidate_sha256 = str(request.get("candidate_sha256") or "").lower()
    slide_index = int(request.get("slide_index") or 0)
    operations = request.get("operations") or []
    if SHA256.fullmatch(expected_candidate_sha256) is None or file_sha256(input_path) != expected_candidate_sha256:
        fail("repair input no longer matches the candidate snapshot", "pptx_visual_repair_stale")
    if not isinstance(operations, list) or not 1 <= len(operations) <= MAX_OPERATIONS:
        fail("repair operation count is outside the allowed range")

    before_package = package_snapshot(input_path)
    presentation = Presentation(input_path)
    if slide_index <= 0 or slide_index > len(presentation.slides):
        fail("repair slide_index is outside the candidate")
    slide = presentation.slides[slide_index - 1]
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)

    bindings = request.get("target_hashes") or {}
    if not isinstance(bindings, dict):
        fail("repair target bindings are invalid")
    referenced = set()
    for operation in operations:
        if not isinstance(operation, dict):
            fail("repair operation must be an object")
        referenced.add(str(operation.get("shape_ref") or ""))
        if operation.get("op") in ("place_above", "place_below"):
            referenced.add(str(operation.get("relative_shape_ref") or ""))
    for shape_ref in referenced:
        shape_index = parse_shape_ref(shape_ref, slide_index)
        shape = require_shape(slide, shape_index)
        expected = str(bindings.get(shape_ref) or "").lower()
        actual = target_hash(expected_candidate_sha256, shape_ref, shape)
        if SHA256.fullmatch(expected) is None or actual != expected:
            fail("repair shape binding is stale for %s (expected %s, actual %s)" % (shape_ref, expected, actual), "pptx_visual_repair_stale")

    deleted = set()
    changed_indexes = set()
    for operation in operations:
        op = str(operation.get("op") or "")
        shape_ref = str(operation.get("shape_ref") or "")
        shape_index = parse_shape_ref(shape_ref, slide_index)
        if shape_index in deleted:
            fail("repair operation targets a shape after deletion")
        shape = require_shape(slide, shape_index)
        if op == "set_geometry":
            if set(operation) != {"op", "shape_ref", "region_milli"}:
                fail("set_geometry contains unsupported fields")
            apply_geometry(shape, operation["region_milli"], slide_width, slide_height)
        elif op == "set_text_style":
            apply_text_style(shape, operation, slide_width, slide_height)
        elif op == "set_shape_style":
            apply_shape_style(shape, operation)
        elif op == "rewrite_text":
            if set(operation) != {"op", "shape_ref", "text"}:
                fail("rewrite_text contains unsupported fields")
            apply_rewrite(shape, operation)
        elif op in ("place_above", "place_below"):
            if set(operation) != {"op", "shape_ref", "relative_shape_ref"}:
                fail("ordering operation contains unsupported fields")
            relative_index = parse_shape_ref(operation["relative_shape_ref"], slide_index)
            if relative_index == shape_index or relative_index in deleted:
                fail("ordering operation has an invalid relative shape")
            relative = require_shape(slide, relative_index)
            tree = shape._element.getparent()
            if tree is None or relative._element.getparent() is not tree:
                fail("ordering operation requires peer shapes")
            tree.remove(shape._element)
            relative_position = tree.index(relative._element)
            tree.insert(relative_position + (1 if op == "place_above" else 0), shape._element)
        elif op == "delete_generated_shape":
            if set(operation) != {"op", "shape_ref", "generated"} or operation.get("generated") is not True:
                fail("delete_generated_shape lacks Runtime generation proof")
            parent = shape._element.getparent()
            if parent is None:
                fail("delete_generated_shape has no parent")
            parent.remove(shape._element)
            deleted.add(shape_index)
        else:
            fail("repair operation is unsupported")
        changed_indexes.add(shape_index)

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    presentation.save(output_path)
    Presentation(output_path)
    after_package = package_snapshot(output_path)
    allowed_part = "ppt/slides/slide%d.xml" % slide_index
    if set(before_package) != set(after_package):
        fail("repair changed the PPTX package part set", "pptx_visual_repair_preservation")
    unexpected = [name for name in before_package if name != allowed_part and before_package[name] != after_package[name]]
    if unexpected:
        fail("repair changed unauthorized package parts: %s" % ", ".join(sorted(unexpected)[:8]), "pptx_visual_repair_preservation")
    if before_package.get(allowed_part) == after_package.get(allowed_part):
        fail("repair produced no semantic slide change")

    print(json.dumps({
        "schema_version": "sparkclaw.pptx_visual_repair_result.v1",
        "slide_index": slide_index,
        "operation_count": len(operations),
        "changed_shape_indexes": sorted(changed_indexes),
        "candidate_sha256": file_sha256(output_path),
        "bytes": os.path.getsize(output_path),
    }, ensure_ascii=False))
except Exception as exc:
    fail(exc)
